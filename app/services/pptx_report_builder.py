"""
PPTX Report Builder - Slate/Electric Blue A4 landscape SPI/SPC analysis report.

Slide structure — SMT SPI Diagnosis-first framework:
  P1.  製程狀態 / Process Status            ← process_state banner、spi_narrative、製程風險
  P2.  主要問題 / Core Diagnosis            ← primary_feature、problem_type_label、異常摘要
  P3a. 製程能力 / Process Capability        ← Cpk/Ppk 分析
  P3b. 製程穩定性 / Process Stability       ← SPC/CUSUM/EWMA 結論
  P5.  多訊號關聯診斷 / Multi-Signal Diagnosis ← 異常訊號來源圖表、關聯分析、製程型態判定
  P6.  製程原因推論 / Process Cause Hypothesis ← 製程原因推論、建議檢查項目
  P3c. 製程風險 / Process Risk              ← 複合風險判斷句
  P4.  Chart Evidence Gallery (optional)   ← 圖表證據
  P5d. Anomaly Diagnosis detail pages      ← 逐條診斷詳頁
  P6d. Distribution Analysis               ← 分布分析證據
  P7.  Spatial Analysis                    ← 空間分析
  P8.  Variation Source Analysis           ← 變異來源
  P9.  Background Info                     ← 工單 + 規格（合併）
  P10. Statistics Summary                  ← 完整統計表
  P11. Appendix                            ← 附錄

Style: light Slate background, Electric Blue accents, color-coded key parameters,
       Traditional Chinese with English professional terms retained.
"""
import io
import logging
from datetime import datetime
from typing import Any, Callable, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.presentation import Presentation as PptxPresentation
from pptx.util import Pt, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from app.analytics.dashboard_layers_display import (
    extract_dashboard_layers,
    pptx_alarm_health_lines,
    pptx_defect_structure_lines,
    pptx_engineering_data_lines,
    pptx_kpi_capability_lines,
    pptx_layer8_diagnosis_lines,
    pptx_stability_dashboard_lines,
    top_refdes_line,
)
from app.services import report_risk
from app.services import report_exec_summary as _rpt_exec
from app.services import multi_signal_diagnosis as _msd
from app.services.diagnostic_evidence_matrix import build_readable_diagnostic_tabs
from app.ui.theme.tokens import (
    ACCENT_ERROR,
    ACCENT_PRIMARY,
    ACCENT_SUCCESS,
    ACCENT_WARNING,
    BG_BLOCK,
    BG_SECONDARY,
    BORDER,
    CHART_PALETTE_AREA_FILL,
    CHART_PALETTE_HEIGHT_FILL,
    CHART_PALETTE_OFFSET_X_FILL,
    CHART_PALETTE_SOLDER_FILL,
    PROCESS_ALARM_CARD_BG_WARNING,
    RPT_BADGE_ERROR_BG,
    RPT_BADGE_INFO_BG,
    RPT_BADGE_NEUTRAL_BG,
    RPT_BADGE_SUCCESS_BG,
    RPT_BADGE_WARNING_BG,
    RPT_SURFACE,
    RPT_TABLE_HEADER_BG,
    TEXT_DISABLED,
    TEXT_PRIMARY,
    TEXT_SECONDARY,
)
from app.utils.numeric_utils import safe_float, coerce_int

logger = logging.getLogger(__name__)

# ── A4 Landscape dimensions ─────────────────────────────────────────────────
SLIDE_WIDTH = Mm(297)
SLIDE_HEIGHT = Mm(210)

# ── Color palette (token-derived, shared with desktop UI) ───────────────────
def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


CLR_BLACK = _rgb(TEXT_PRIMARY)
CLR_TITLE = _rgb(TEXT_PRIMARY)
CLR_SUBTITLE = _rgb(TEXT_SECONDARY)
CLR_WHITE = _rgb(BG_BLOCK)
CLR_GOOD = _rgb(ACCENT_SUCCESS)
CLR_WARNING = _rgb(ACCENT_WARNING)
CLR_BAD = _rgb(ACCENT_ERROR)
CLR_ACCENT = _rgb(ACCENT_PRIMARY)
CLR_LIGHT_BG = _rgb(RPT_SURFACE)
CLR_BORDER = _rgb(BORDER)
CLR_TABLE_HEADER_BG = _rgb(RPT_TABLE_HEADER_BG)
CLR_TABLE_HEADER_TEXT = _rgb(TEXT_PRIMARY)
CLR_TABLE_ALT_BG = _rgb(BG_SECONDARY)
CLR_FOOTER = _rgb(TEXT_DISABLED)
CLR_INDICATOR_BG = _rgb(BG_BLOCK)
CLR_BADGE_INFO_BG = _rgb(RPT_BADGE_INFO_BG)
CLR_BADGE_INFO_BORDER = _rgb(ACCENT_PRIMARY)
CLR_BADGE_INFO_TEXT = _rgb(ACCENT_PRIMARY)
CLR_BADGE_WARNING_BG = _rgb(RPT_BADGE_WARNING_BG)
CLR_BADGE_WARNING_BORDER = _rgb(ACCENT_WARNING)
CLR_BADGE_WARNING_TEXT = _rgb(ACCENT_WARNING)
CLR_BADGE_ERROR_BG = _rgb(RPT_BADGE_ERROR_BG)
CLR_BADGE_SUCCESS_BG = _rgb(RPT_BADGE_SUCCESS_BG)
CLR_BADGE_NEUTRAL_BG = _rgb(RPT_BADGE_NEUTRAL_BG)

# ── Font ─────────────────────────────────────────────────────────────────────
FONT_FAMILY = "Noto Sans TC"
FONT_FAMILY_FALLBACK = "Microsoft JhengHei"

# ── Margins ──────────────────────────────────────────────────────────────────
MARGIN_LEFT = Mm(16)
MARGIN_TOP = Mm(14)
MARGIN_RIGHT = Mm(16)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT  # ~265mm


def _font(run, size_pt=11, bold=False, color=None, font_name=None):
    """Apply font styling to a run."""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color or CLR_BLACK
    run.font.name = font_name or FONT_FAMILY


def _add_textbox(slide, left, top, width, height, text="", size_pt=11,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name=None):
    """Add a styled text box and return the textbox shape."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _font(run, size_pt, bold, color, font_name)
    return txbox


def _add_list_box(slide, left, top, width, height, lines: List[str], size_pt=9, color=None):
    """Add a multi-line text box, one bullet-style line per paragraph."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.clear()
    normalized = [str(line).strip() for line in lines if str(line).strip()]
    if not normalized:
        normalized = ["—"]
    for index, line in enumerate(normalized):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(3)
        run = paragraph.add_run()
        run.text = f"- {line}"
        _font(run, size_pt, False, color)
    return txbox


def _add_colored_list_box(
    slide,
    left,
    top,
    width,
    height,
    lines: List[str],
    *,
    size_pt: int = 9,
    default_color: Optional[RGBColor] = None,
    color_resolver: Optional[Callable[[str], RGBColor]] = None,
):
    """Add bullet lines with per-line color resolution."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.clear()
    normalized = [str(line).strip() for line in lines if str(line).strip()]
    if not normalized:
        normalized = ["—"]
    for index, line in enumerate(normalized):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(3)
        run = paragraph.add_run()
        run.text = f"- {line}"
        line_color = color_resolver(line) if color_resolver else (default_color or CLR_BLACK)
        _font(run, size_pt, False, line_color)
    return txbox


def _ellipsize(text: Any, max_chars: int) -> str:
    """Trim overly long text for stable A4 slide layout."""
    normalized = " ".join(str(text or "").split())
    if len(normalized) <= max_chars:
        return normalized
    return normalized[: max_chars - 1].rstrip() + "…"


def _add_footer(slide, timestamp_str: str, page_num: int, total_pages: int):
    """Add a footer bar with timestamp and page number."""
    footer_top = SLIDE_HEIGHT - Mm(12)
    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, footer_top - Mm(1),
        CONTENT_WIDTH, Mm(0.4),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_BORDER
    line.line.fill.background()

    # Left: timestamp
    _add_textbox(
        slide, MARGIN_LEFT, footer_top, Mm(120), Mm(10),
        text=f"產生時間：{timestamp_str}",
        size_pt=7, color=CLR_FOOTER,
    )
    # Right: page number
    _add_textbox(
        slide, SLIDE_WIDTH - MARGIN_RIGHT - Mm(40), footer_top, Mm(40), Mm(10),
        text=f"{page_num} / {total_pages}",
        size_pt=7, color=CLR_FOOTER, alignment=PP_ALIGN.RIGHT,
    )


def _add_slide_title(slide, title_text: str, subtitle_text: str = ""):
    """Add a slide title bar with accent underline."""
    _add_textbox(
        slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Mm(12),
        text=title_text, size_pt=22, bold=True, color=CLR_TITLE,
    )
    # Accent underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, MARGIN_TOP + Mm(13),
        Mm(48), Mm(1),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_ACCENT
    line.line.fill.background()

    if subtitle_text:
        _add_textbox(
            slide, MARGIN_LEFT, MARGIN_TOP + Mm(16), CONTENT_WIDTH, Mm(8),
            text=subtitle_text, size_pt=10, color=CLR_SUBTITLE,
        )


def _set_cell_text(cell, text: str, size_pt=9, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Set styled text in a table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    _font(run, size_pt, bold, color)


def _style_table_header(table, col_count: int):
    """Style the header row of a table."""
    for ci in range(col_count):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_TABLE_HEADER_BG
        # Add bottom border look (by coloring cell or just text)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(9)
                run.font.color.rgb = CLR_TABLE_HEADER_TEXT


def _style_table_rows(table, row_start: int, row_end: int, col_count: int):
    """Apply alternating row backgrounds."""
    for ri in range(row_start, row_end):
        if ri % 2 == 1:
            for ci in range(col_count):
                cell = table.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_TABLE_ALT_BG


def _add_status_indicator(slide, left, top, width, height, label, value, color, icon_text=""):
    """Add a compact status indicator box."""
    # Outer box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    box.fill.solid()
    box.fill.fore_color.rgb = CLR_INDICATOR_BG
    box.line.color.rgb = CLR_BORDER
    box.line.width = Pt(0.5)

    # Indicator light (small circle)
    circle_size = Mm(4)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Mm(3), top + Mm(3), circle_size, circle_size,
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    # Label
    _add_textbox(
        slide, left + Mm(8), top + Mm(2.5), width - Mm(10), Mm(5),
        text=label, size_pt=7, bold=True, color=CLR_SUBTITLE,
    )
    # Value
    _add_textbox(
        slide, left + Mm(4), top + Mm(8), width - Mm(8), Mm(10),
        text=value, size_pt=14, bold=True, color=color, alignment=PP_ALIGN.CENTER,
    )


def _cpk_color(value) -> RGBColor:
    """Return color for Cpk/Ppk values."""
    if value is None:
        return CLR_BLACK
    try:
        v = float(str(value).split()[0])
    except (ValueError, IndexError):
        return CLR_BLACK
    if v >= 1.67:
        return CLR_GOOD
    if v >= 1.33:
        return CLR_WARNING
    return CLR_BAD


def _yield_color(value) -> RGBColor:
    """Return color for yield percentage."""
    if value is None:
        return CLR_BLACK
    try:
        v = float(value)
    except (ValueError, TypeError):
        return CLR_BLACK
    if 0.0 <= v <= 1.0:
        v *= 100.0
    if v >= 99.0:
        return CLR_GOOD
    if v >= 95.0:
        return CLR_WARNING
    return CLR_BAD


def _verdict_color(text: str) -> RGBColor:
    """Return color for verdict text."""
    if "可接受" in text:
        return CLR_GOOD
    if "不可接受" in text or "拒收" in text:
        return CLR_BAD
    if "待改善" in text:
        return CLR_WARNING
    return CLR_BLACK


def _severity_color(severity: str) -> RGBColor:
    """Map root-cause severity to color."""
    if severity == "error":
        return CLR_BAD
    if severity == "warning":
        return CLR_WARNING
    return CLR_ACCENT


def _severity_label(severity: str) -> str:
    """Map severity to Chinese label."""
    return {"error": "嚴重", "warning": "警示", "info": "提示"}.get(severity, "提示")


def _first_numeric_token(text: str) -> Optional[float]:
    """Extract first numeric token from a line such as 'Cpk: 0.92'."""
    import re

    normalized = str(text or "").replace(",", "")
    match = re.search(r"-?\d+(?:\.\d+)?", normalized)
    if not match:
        return None
    try:
        return float(match.group(0))
    except (TypeError, ValueError):
        return None


def _evidence_line_color(line: str) -> RGBColor:
    """Highlight key evidence metrics while keeping default black text."""
    text = str(line or "").strip()
    if not text:
        return CLR_BLACK
    lower = text.lower()
    numeric = _first_numeric_token(text)

    if "cpk" in lower or lower.startswith("cp:") or lower.startswith("cp "):
        return _cpk_color(numeric)

    if "yield" in lower or "良率" in text:
        return _yield_color(numeric)

    if "ppm" in lower or "dpmo" in lower:
        if numeric is None:
            return CLR_ACCENT
        if numeric <= 100:
            return CLR_GOOD
        if numeric <= 10000:
            return CLR_WARNING
        return CLR_BAD

    if (
        "ooc ratio" in lower
        or "oos ratio" in lower
        or "ooc 比率" in lower
        or "oos 比率" in lower
        or "decline ratio" in lower
        or "edge oos ratio" in lower
    ):
        if numeric is None:
            return CLR_ACCENT
        ratio = numeric / 100.0 if "%" in text else numeric
        if "oos" in lower or "OOS" in text:
            return CLR_BAD if ratio > 0 else CLR_GOOD
        if ratio >= 0.10:
            return CLR_BAD
        if ratio >= 0.03:
            return CLR_WARNING
        return CLR_GOOD

    if "variance ratio" in lower:
        if numeric is None:
            return CLR_ACCENT
        if numeric >= 2.0:
            return CLR_BAD
        if numeric >= 1.3:
            return CLR_WARNING
        return CLR_GOOD

    if "p-value" in lower:
        if numeric is None:
            return CLR_ACCENT
        return CLR_BAD if numeric < 0.05 else CLR_GOOD

    if "is normal" in lower:
        if "false" in lower:
            return CLR_BAD
        if "true" in lower:
            return CLR_GOOD
        return CLR_ACCENT

    if any(keyword in lower for keyword in ("error", "warning", "異常", "超規", "失敗", "嚴重")):
        return CLR_BAD

    return CLR_BLACK


def _add_picture_contain(slide, chart_bytes: bytes, left, top, width, height):
    """Render an image into a bounding box while keeping aspect ratio."""
    img_stream = io.BytesIO(chart_bytes)
    try:
        from PIL import Image

        with Image.open(io.BytesIO(chart_bytes)) as image:
            img_w, img_h = image.size
        if img_w <= 0 or img_h <= 0:
            raise ValueError("invalid image size")
        box_ratio = float(width) / float(height)
        img_ratio = float(img_w) / float(img_h)
        if img_ratio >= box_ratio:
            draw_w = int(width)
            draw_h = int(draw_w / img_ratio)
        else:
            draw_h = int(height)
            draw_w = int(draw_h * img_ratio)
        draw_left = int(left + (width - draw_w) / 2)
        draw_top = int(top + (height - draw_h) / 2)
        slide.shapes.add_picture(img_stream, draw_left, draw_top, draw_w, draw_h)
    except (ImportError, OSError, TypeError, ValueError, RuntimeError):
        logger.debug("PPTX: fallback to stretched image placement", exc_info=True)
        slide.shapes.add_picture(img_stream, left, top, width, height)


def _build_slide_bullet_section(
    prs: PptxPresentation,
    *,
    title_text: str,
    subtitle_text: str,
    lines: List[str],
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    size_pt: int = 10,
    colored: bool = False,
):
    """Build a generic section slide with one content card and bullet lines."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, title_text, subtitle_text)
    content_top = MARGIN_TOP + Mm(28)
    content_h = Mm(130)

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT, content_top, CONTENT_WIDTH, content_h,
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CLR_LIGHT_BG
    card.line.color.rgb = CLR_BORDER
    card.line.width = Pt(0.5)

    normalized = [_ellipsize(line, 110) for line in lines if str(line).strip()]
    if not normalized:
        normalized = ["資料不足，請確認來源資料與分析條件。"]
    if colored:
        _add_colored_list_box(
            slide,
            MARGIN_LEFT + Mm(6),
            content_top + Mm(7),
            CONTENT_WIDTH - Mm(12),
            content_h - Mm(12),
            normalized,
            size_pt=size_pt,
            default_color=CLR_BLACK,
            color_resolver=_evidence_line_color,
        )
    else:
        _add_list_box(
            slide,
            MARGIN_LEFT + Mm(6),
            content_top + Mm(7),
            CONTENT_WIDTH - Mm(12),
            content_h - Mm(12),
            normalized,
            size_pt=size_pt,
            color=CLR_BLACK,
        )

    _add_footer(slide, timestamp_str, page_num, total_pages)


def _find_diagnostic_by_keywords(diagnostics: List[Dict[str, Any]], keywords: List[str]) -> Optional[Dict[str, Any]]:
    """Return first diagnostic entry containing any keyword in rule/title/summary/charts."""
    normalized_keywords = [k.lower() for k in keywords if str(k).strip()]
    if not normalized_keywords:
        return None
    for item in diagnostics:
        chart_tokens = " ".join(item.get("observable_charts", []) or [])
        candidate = " ".join(
            [
                str(item.get("rule_id", "")),
                str(item.get("chart_title", "")),
                str(item.get("summary", "")),
                chart_tokens,
            ]
        ).lower()
        if any(token in candidate for token in normalized_keywords):
            return item
    return None


def _dedupe_chart_ids(chart_ids: List[str]) -> List[str]:
    """Keep chart order stable while removing duplicates/empties."""
    ordered: List[str] = []
    seen: set[str] = set()
    for chart_id in chart_ids:
        cid = str(chart_id or "").strip()
        if not cid or cid in seen:
            continue
        seen.add(cid)
        ordered.append(cid)
    return ordered


def _get_chart_required_feature_count(chart_id: str) -> int:
    """Resolve chart feature-count contract from chart registry."""
    from app.analytics.chart_registry import CHART_CATALOG

    for entry in CHART_CATALOG:
        if str(entry.get("id", "")) == chart_id:
            raw_count = entry.get("required_feature_count", 1)
            if isinstance(raw_count, bool):
                return 1
            if isinstance(raw_count, (int, float)):
                return int(raw_count)
            if isinstance(raw_count, str):
                try:
                    return int(raw_count)
                except ValueError:
                    return 1
            return 1
    return 1


def _resolve_chart_features_for_export(
    chart_id: str,
    *,
    selected_features: List[str],
    available_features: List[str],
) -> List[str]:
    """Pick feature list for a chart render according to its required feature count."""
    ordered: List[str] = []
    for feature in [*selected_features, *available_features]:
        normalized = str(feature or "").strip()
        if not normalized or normalized in ordered:
            continue
        ordered.append(normalized)
    required = _get_chart_required_feature_count(chart_id)
    if required <= 1:
        return ordered[:1]
    return ordered[:required]


def _render_chart_evidence_items(
    *,
    selected_chart_ids: List[str],
    analysis_payload: Dict[str, Any],
    selected_features: List[str],
    available_features: List[str],
    coverage_by_id: Optional[Dict[str, Dict[str, Any]]] = None,
    render_chart_fn: Optional[Callable[..., Optional[bytes]]] = None,
) -> List[Dict[str, Any]]:
    """Render selected charts into image evidence items for PPTX gallery slides."""
    if not selected_chart_ids or not isinstance(analysis_payload, dict) or not analysis_payload:
        return []

    from app.analytics.chart_registry import get_chart_display_name
    if render_chart_fn is None:
        from app.services.chart_render import render_chart_to_png_bytes
        render_chart_fn = render_chart_to_png_bytes

    items: List[Dict[str, Any]] = []
    for chart_id in _dedupe_chart_ids(selected_chart_ids):
        coverage_item = (coverage_by_id or {}).get(chart_id)
        if coverage_item is not None and coverage_item.get("status") == "未納入":
            continue
        required = _get_chart_required_feature_count(chart_id)
        features = _resolve_chart_features_for_export(
            chart_id,
            selected_features=selected_features,
            available_features=available_features,
        )
        if len(features) < required:
            if coverage_item is not None:
                coverage_item["status"] = "不相容"
                coverage_item["reason"] = f"需 {required} 特徵"
            continue
        try:
            chart_bytes = render_chart_fn(
                chart_id,
                analysis_payload,
                features=features,
                context="report",
            )
        except (AttributeError, KeyError, TypeError, ValueError, RuntimeError, OSError):
            logger.exception("PPTX: 圖表證據渲染失敗: chart_id=%s", chart_id)
            if coverage_item is not None:
                coverage_item["status"] = "渲染失敗"
                coverage_item["reason"] = "圖表渲染失敗"
            continue
        if not chart_bytes:
            if coverage_item is not None:
                coverage_item["status"] = "無資料"
                coverage_item["reason"] = "圖表資料不足或未產生有效圖像"
            continue
        if coverage_item is not None:
            coverage_item["status"] = "已輸出"
            coverage_item["features"] = features
            if len(features) == 1 and len(selected_features) > 1:
                coverage_item["reason"] = f"已輸出；代表特徵：{features[0]}"
            else:
                coverage_item["reason"] = "已納入 PPTX 圖表證據畫廊"
        items.append(
            {
                "chart_id": chart_id,
                "title": get_chart_display_name(chart_id, lang="zh_only"),
                "features": features,
                "image_bytes": chart_bytes,
            }
        )
    return items


def _build_slide_chart_evidence_gallery(
    prs: PptxPresentation,
    *,
    chart_items: List[Dict[str, Any]],
    gallery_index: int,
    gallery_total: int,
    timestamp_str: str,
    page_num: int,
    total_pages: int,
) -> None:
    """Build a 2x2 chart evidence gallery slide for selected report charts."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    suffix = f" ({gallery_index}/{gallery_total})" if gallery_total > 1 else ""
    _add_slide_title(
        slide,
        f"5A. Chart Evidence Gallery / 數據與圖表證據{suffix}",
        "數據特徵診斷證據頁面 (Process Analytics Evidence)",
    )

    grid_top = MARGIN_TOP + Mm(28)
    grid_bottom = SLIDE_HEIGHT - Mm(20)
    grid_height = int(grid_bottom - grid_top)
    gap_x = int(Mm(6))
    gap_y = int(Mm(6))
    cell_w = int((CONTENT_WIDTH - gap_x) / 2)
    cell_h = int((grid_height - gap_y) / 2)

    for idx, item in enumerate(chart_items[:4]):
        row = idx // 2
        col = idx % 2
        left = int(MARGIN_LEFT + col * (cell_w + gap_x))
        top = int(grid_top + row * (cell_h + gap_y))

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, cell_w, cell_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CLR_LIGHT_BG
        card.line.color.rgb = CLR_BORDER
        card.line.width = Pt(0.5)

        chart_title = _ellipsize(item.get("title", "Chart"), 36)
        feature_line = ", ".join(item.get("features", [])[:3]) or "—"
        _add_textbox(
            slide,
            left + Mm(3),
            top + Mm(2.5),
            cell_w - Mm(6),
            Mm(8),
            text=chart_title,
            size_pt=9,
            bold=True,
            color=CLR_TITLE,
        )
        _add_textbox(
            slide,
            left + Mm(3),
            top + Mm(10),
            cell_w - Mm(6),
            Mm(6),
            text=f"特徵：{_ellipsize(feature_line, 42)}",
            size_pt=7,
            color=CLR_SUBTITLE,
        )

        image_left = int(left + Mm(3))
        image_top = int(top + Mm(16.5))
        image_w = int(cell_w - Mm(6))
        image_h = int(cell_h - Mm(19.5))
        _add_picture_contain(slide, item["image_bytes"], image_left, image_top, image_w, image_h)

    _add_footer(slide, timestamp_str, page_num, total_pages)


def _is_missing_value(value: Any) -> bool:
    """Return True when a text-like value should be treated as missing."""
    text = str(value or "").strip()
    return text in {"", "—", "N/A", "UNKNOWN", "VERIFY", "None"}


def _report_value(value: Any, *, unknown_text: str = "UNKNOWN (VERIFY)") -> str:
    """Normalize display values so missing fields are explicit in report."""
    return unknown_text if _is_missing_value(value) else str(value).strip()


def _format_optional_mm(value: Any) -> str:
    """Format optional numeric thickness in mm."""
    parsed = _to_float(value)
    if parsed is None:
        return "—"
    return f"{parsed:.3f} mm"


def _collect_filter_context_lines(report_context: Dict[str, Any]) -> List[str]:
    """Build report scope lines from filter context."""
    filters = report_context.get("filter_context", {}) if isinstance(report_context, dict) else {}
    if not isinstance(filters, dict):
        return []
    label_map = {
        "batch": "Batch",
        "refdes": "RefDes",
        "part_type": "PartType",
        "product": "Product",
        "time_start": "Time Start",
        "time_end": "Time End",
        "line": "Line",
    }
    lines: List[str] = []
    for key in ("batch", "refdes", "part_type", "product", "time_start", "time_end", "line"):
        value = filters.get(key)
        if _is_missing_value(value):
            continue
        text = str(value).strip()
        if text in {"全部 (All)", "全部"}:
            continue
        lines.append(f"{label_map.get(key, key)}: {text}")
    return lines


def _coverage_items(report_context: Dict[str, Any]) -> List[Dict[str, Any]]:
    coverage = report_context.get("evidence_coverage", {}) if isinstance(report_context, dict) else {}
    items = coverage.get("items", []) if isinstance(coverage, dict) else []
    return [item for item in items if isinstance(item, dict)]


def _coverage_item_by_id(report_context: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    return {
        str(item.get("chart_id", "")).strip(): item
        for item in _coverage_items(report_context)
        if str(item.get("chart_id", "")).strip()
    }


def _data_scope_lines(report_context: Dict[str, Any]) -> List[str]:
    scope = report_context.get("data_scope", {}) if isinstance(report_context, dict) else {}
    if not isinstance(scope, dict):
        return []
    used_sources = [str(item) for item in (scope.get("used_sources") or []) if str(item).strip()]
    selected_features = [
        str(item) for item in (scope.get("selected_features") or []) if str(item).strip()
    ]
    lines = [
        f"本次使用資料：{', '.join(used_sources) if used_sources else '量測資料'}",
        f"樣本數：{int(scope.get('sample_n') or 0):,}",
        f"分析特徵：{', '.join(selected_features) if selected_features else '—'}",
    ]
    filter_lines = _collect_filter_context_lines(report_context)
    if filter_lines:
        lines.append(f"篩選條件：{'; '.join(filter_lines)}")
    excluded = scope.get("excluded_evidence") or []
    excluded_labels = [
        f"{item.get('label', '—')}（{item.get('reason', '—')}）"
        for item in excluded
        if isinstance(item, dict)
    ]
    if excluded_labels:
        lines.append(f"本次未納入：{'; '.join(excluded_labels)}")
    return lines


def _trust_status_lines(report_context: Dict[str, Any]) -> List[str]:
    scope = report_context.get("data_scope", {}) if isinstance(report_context, dict) else {}
    statuses = scope.get("section_trust", {}) if isinstance(scope, dict) else {}
    if not isinstance(statuses, dict):
        return []
    return [
        f"統計/能力：{statuses.get('statistics', '可信：資料直接計算')}",
        f"圖表證據：{statuses.get('charts', '可信：圖表證據')}",
        f"根因/對策：{statuses.get('inference', '需複核：規則推論')}",
        f"空間分析：{statuses.get('spatial', '未納入：資料缺失')}",
    ]


def _has_coordinate_scope(report_context: Dict[str, Any], spatial_payload: Dict[str, Any]) -> bool:
    scope = report_context.get("data_scope", {}) if isinstance(report_context, dict) else {}
    if isinstance(scope, dict) and "has_coordinate_data" in scope:
        return bool(scope.get("has_coordinate_data"))
    # Backward-compatible direct builder calls may not provide data_scope.
    if isinstance(spatial_payload, dict) and spatial_payload.get("statistics"):
        return True
    relation_meta = report_context.get("relation_meta", {}) if isinstance(report_context, dict) else {}
    return bool(isinstance(relation_meta, dict) and relation_meta.get("match_rate") is not None)


def _coverage_pages_count(report_context: Dict[str, Any], rows_per_page: int = 14) -> int:
    items = _coverage_items(report_context)
    if not items:
        return 0
    return (len(items) + rows_per_page - 1) // rows_per_page


def _build_slide_chart_evidence_coverage(
    prs: PptxPresentation,
    *,
    report_context: Dict[str, Any],
    page_index: int,
    page_total: int,
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    rows_per_page: int = 14,
) -> None:
    items = _coverage_items(report_context)
    start = (page_index - 1) * rows_per_page
    page_items = items[start : start + rows_per_page]
    suffix = f" ({page_index}/{page_total})" if page_total > 1 else ""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(
        slide,
        f"圖表證據覆蓋表 / Evidence Coverage{suffix}",
        "圖表名稱 · 使用特徵 · 狀態 · 原因",
    )
    content_top = MARGIN_TOP + Mm(28)
    summary = (report_context.get("evidence_coverage", {}) or {}).get("summary", {})
    if isinstance(summary, dict):
        summary_line = (
            f"總數 {summary.get('total', len(items))}；"
            f"可用 {summary.get('available', 0)}；"
            f"不相容 {summary.get('incompatible', 0)}；"
            f"未納入 {summary.get('excluded', 0)}；"
            f"勾選 {summary.get('selected', 0)}"
        )
        _add_textbox(
            slide,
            MARGIN_LEFT,
            content_top,
            CONTENT_WIDTH,
            Mm(6),
            text=summary_line,
            size_pt=8,
            color=CLR_SUBTITLE,
        )
        table_top = content_top + Mm(8)
    else:
        table_top = content_top

    row_count = max(2, len(page_items) + 1)
    table = slide.shapes.add_table(
        row_count,
        4,
        MARGIN_LEFT,
        table_top,
        CONTENT_WIDTH,
        Mm(112),
    ).table
    widths = [Mm(72), Mm(58), Mm(30), CONTENT_WIDTH - Mm(160)]
    for idx, width in enumerate(widths):
        table.columns[idx].width = int(width)
    _style_table_header(table, 4)
    _set_cell_text(table.cell(0, 0), "圖表", size_pt=7, bold=True)
    _set_cell_text(table.cell(0, 1), "使用特徵", size_pt=7, bold=True)
    _set_cell_text(table.cell(0, 2), "狀態", size_pt=7, bold=True, alignment=PP_ALIGN.CENTER)
    _set_cell_text(table.cell(0, 3), "原因", size_pt=7, bold=True)
    for row, item in enumerate(page_items, start=1):
        features = item.get("features") or []
        feature_text = " + ".join(str(f) for f in features if str(f).strip()) or "—"
        status = str(item.get("status") or "—")
        reason = str(item.get("reason") or "")
        _set_cell_text(table.cell(row, 0), _ellipsize(item.get("chart_name", "—"), 24), size_pt=6)
        _set_cell_text(table.cell(row, 1), _ellipsize(feature_text, 28), size_pt=6)
        _set_cell_text(
            table.cell(row, 2),
            status,
            size_pt=6,
            bold=status in {"已輸出", "未納入", "不相容", "渲染失敗"},
            color=(
                CLR_GOOD if status == "已輸出"
                else CLR_WARNING if status in {"未納入", "不相容", "無資料"}
                else CLR_BAD if status == "渲染失敗"
                else CLR_SUBTITLE
            ),
            alignment=PP_ALIGN.CENTER,
        )
        _set_cell_text(table.cell(row, 3), _ellipsize(reason or "—", 38), size_pt=6)
    _style_table_rows(table, 1, row_count, 4)
    _add_footer(slide, timestamp_str, page_num, total_pages)


def _resolve_workorder_field(
    wo_master: Dict[str, Any],
    report_context: Dict[str, Any],
    *,
    key: str,
    fallback_key: Optional[str] = None,
) -> Any:
    """Resolve field from workorder master with optional coordinate-registry fallback."""
    value = wo_master.get(key, "—")
    if not _is_missing_value(value):
        return value
    entry = report_context.get("coordinate_registry_entry", {}) if isinstance(report_context, dict) else {}
    if isinstance(entry, dict) and fallback_key:
        fallback = entry.get(fallback_key)
        if not _is_missing_value(fallback):
            return fallback
    return value


# ═══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════════════════════

def _format_process_diagnosis_report_lines(pdr: Dict[str, Any]) -> List[str]:
    """Turn process_diagnosis_report DTO into bullet lines for the four-layer slide."""
    _ad = pdr.get("A_decision")
    _bd = pdr.get("B_diagnosis")
    _cd = pdr.get("C_evidence")
    _dd = pdr.get("D_data")
    a: Dict[str, Any] = _ad if isinstance(_ad, dict) else {}
    b: Dict[str, Any] = _bd if isinstance(_bd, dict) else {}
    c: Dict[str, Any] = _cd if isinstance(_cd, dict) else {}
    d: Dict[str, Any] = _dd if isinstance(_dd, dict) else {}
    pats = b.get("process_patterns") or []
    pat_s = "、".join(str(p) for p in pats[:8]) if isinstance(pats, list) and pats else "—"
    core = _ellipsize(str(a.get("core_diagnosis_zh") or ""), 200)
    pl = str(a.get("pattern_logic") or "").strip()
    pl_line = f"  模式邏輯：{_ellipsize(pl, 180)}" if pl else ""
    lines = [
        f"[A 決策層] 判定 {a.get('process_verdict', '—')} · 風險 {a.get('risk_level', '—')}",
        f"  {core}",
    ]
    if pl_line:
        lines.append(pl_line)
    lines.extend(
        [
            f"[B 診斷層] 範圍 {b.get('scope', '—')} · 分布 {b.get('distribution_shape', '—')}"
            f" · 假設域 {b.get('hypothesis_domain', '—')}",
            f"  製程型態：{pat_s}",
            f"[C 證據層] {_ellipsize(str(c.get('bridge_zh') or '—'), 220)}",
            f"  組合覆蓋：{c.get('combination_coverage_zh', '—')}",
            f"[D 資料層] {_ellipsize(str(d.get('bridge_zh') or '—'), 220)}",
        ]
    )
    top = c.get("top_evidence")
    readable = c.get("readable_evidence")
    if isinstance(readable, list) and readable:
        for item in readable[:3]:
            if not isinstance(item, dict):
                continue
            lines.append(
                "  證據："
                f"{item.get('title', '—')} / {item.get('result_zh', '—')} / "
                f"{_ellipsize(str(item.get('reason_zh') or '—'), 92)}"
            )
        return lines
    if isinstance(top, list) and top:
        for item in top[:3]:
            if not isinstance(item, dict):
                continue
            feature_s = " + ".join(str(x) for x in (item.get("feature_set") or []))
            lines.append(
                "  證據："
                f"{item.get('chart_name', '—')} / {feature_s or '—'} / "
                f"{_ellipsize(str(item.get('metric_snapshot') or '—'), 80)}"
            )
    return lines


# ── Exec-summary data helper ─────────────────────────────────────────────────

def _build_exec_summary_data(
    summary_data: Dict[str, Any],
    diagnostics: List[Dict[str, Any]],
    risk_assessment: Dict[str, Any],
    *,
    report_context: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """Delegate to report_exec_summary to produce the structured exec-summary dict."""
    rc = report_context if isinstance(report_context, dict) else {}
    dn = rc.get("decision_narrative")
    decision_narrative = dn if isinstance(dn, dict) else None
    return _rpt_exec.build_executive_summary_pptx_data(
        summary_data,
        diagnostics,
        risk_assessment,
        generate_one_liner_fn=report_risk.generate_one_liner,
        generate_risk_sentence_fn=report_risk.generate_risk_sentence,
        derive_stability_verdict_fn=report_risk.derive_stability_verdict,
        requires_immediate_action_fn=report_risk.requires_immediate_action,
        anomaly_type_label_fn=report_risk.anomaly_type_label,
        # SMT SPI process diagnosis
        derive_process_state_fn=report_risk.derive_process_state,
        derive_problem_type_fn=report_risk.derive_problem_type,
        problem_type_zh_fn=report_risk.problem_type_zh,
        generate_spi_narrative_fn=report_risk.generate_spi_narrative,
        decision_narrative=decision_narrative,
    )


# ── P1: Executive Summary ────────────────────────────────────────────────────

def _build_slide_executive_summary(
    prs: PptxPresentation,
    exec_data: Dict[str, Any],
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    *,
    dashboard_layers: Optional[Dict[str, Any]] = None,
    report_context: Optional[Dict[str, Any]] = None,
):
    """
    P1 — 製程狀態 / Process Status  (SMT SPI engineering dashboard).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "製程狀態 / Process Status", "SMT SPI 製程統計分析彙報")

    content_top = MARGIN_TOP + Mm(26)
    indicator_h = Mm(22)
    indicator_w = (CONTENT_WIDTH - Mm(10)) / 3

    # ── Status Indicators Row ────────────────────────────────────────────────
    # 1. Capability
    min_cpk = exec_data.get("min_cpk_value", 0.0)
    cpk_label = f"能力：{exec_data.get('min_cpk_str', '—')}"
    _add_status_indicator(
        slide, MARGIN_LEFT, content_top, indicator_w, indicator_h,
        "製程能力 (Capability)", cpk_label, _cpk_color(min_cpk)
    )

    # 2. Stability
    stability = str(exec_data.get("stability_verdict", "穩定"))
    stab_clr = (
        CLR_BAD if stability == "失控"
        else CLR_WARNING if stability in ("警示", "漂移", "偏移")
        else CLR_GOOD
    )
    _add_status_indicator(
        slide, MARGIN_LEFT + indicator_w + Mm(5), content_top, indicator_w, indicator_h,
        "穩定性 (Stability)", f"製程{stability}", stab_clr
    )

    # 3. Process Health (Health Score)
    risk_level = str(exec_data.get("risk_level", "LOW")).upper()
    health_clr = {"HIGH": CLR_BAD, "MEDIUM": CLR_WARNING, "LOW": CLR_GOOD}.get(risk_level, CLR_BLACK)
    health_label = {"HIGH": "高風險 (High)", "MEDIUM": "中風險 (Med)", "LOW": "良好 (Low)"}.get(risk_level, risk_level)
    _add_status_indicator(
        slide, MARGIN_LEFT + (indicator_w + Mm(5)) * 2, content_top, indicator_w, indicator_h,
        "製程健康度 (Health)", health_label, health_clr
    )

    y_panel = content_top + indicator_h + Mm(6)
    panel_h = SLIDE_HEIGHT - y_panel - Mm(18)
    left_w = Mm(163)
    right_left = MARGIN_LEFT + left_w + Mm(5)
    right_w = CONTENT_WIDTH - left_w - Mm(5)

    # ── Left panel: Diagnosis & Logic ────────────────────────────────────────
    lp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT, y_panel, left_w, panel_h,
    )
    lp.fill.solid()
    lp.fill.fore_color.rgb = CLR_LIGHT_BG
    lp.line.color.rgb = CLR_BORDER
    lp.line.width = Pt(0.5)

    # Problem Type Badge
    problem_type_label = str(exec_data.get("problem_type_label", ""))
    if problem_type_label and problem_type_label not in ("—", "Unknown", ""):
        _add_textbox(
            slide, MARGIN_LEFT + Mm(5), y_panel + Mm(4), left_w - Mm(10), Mm(8),
            text=f"問題定位 : {problem_type_label}", size_pt=10, bold=True, color=CLR_ACCENT
        )
        diag_y = y_panel + Mm(14)
    else:
        diag_y = y_panel + Mm(5)

    # Narrative
    spi_narrative = str(exec_data.get("spi_narrative", "")).strip() or str(exec_data.get("one_liner", "—"))
    _add_textbox(slide, MARGIN_LEFT + Mm(5), diag_y, Mm(50), Mm(6), text="【工程診斷摘要】", size_pt=8, bold=True, color=CLR_SUBTITLE)
    narr_box = slide.shapes.add_textbox(MARGIN_LEFT + Mm(5), diag_y + Mm(6), left_w - Mm(10), Mm(32))
    narr_box.text_frame.word_wrap = True
    _p = narr_box.text_frame.paragraphs[0]
    _run = _p.add_run()
    _run.text = spi_narrative
    _font(_run, 9, color=CLR_BLACK)

    # Risk Sentence
    rs_y = diag_y + Mm(40)
    _add_textbox(slide, MARGIN_LEFT + Mm(5), rs_y, Mm(50), Mm(6), text="【品質風險評量】", size_pt=8, bold=True, color=CLR_SUBTITLE)
    rs_box = slide.shapes.add_textbox(MARGIN_LEFT + Mm(5), rs_y + Mm(6), left_w - Mm(10), Mm(26))
    rs_box.text_frame.word_wrap = True
    _p = rs_box.text_frame.paragraphs[0]
    _run = _p.add_run()
    _run.text = str(exec_data.get("risk_sentence", "—"))
    _font(_run, 9, color=CLR_BLACK)

    scope_lines = _data_scope_lines(report_context or {})[:3]
    trust_lines = _trust_status_lines(report_context or {})[:2]
    disclosure_lines = scope_lines + trust_lines
    if disclosure_lines:
        ds_y = rs_y + Mm(34)
        _add_textbox(
            slide,
            MARGIN_LEFT + Mm(5),
            ds_y,
            Mm(58),
            Mm(6),
            text="【資料來源與可信度】",
            size_pt=8,
            bold=True,
            color=CLR_SUBTITLE,
        )
        _add_list_box(
            slide,
            MARGIN_LEFT + Mm(5),
            ds_y + Mm(7),
            left_w - Mm(10),
            Mm(34),
            disclosure_lines[:5],
            size_pt=6,
            color=CLR_BLACK,
        )

    # ── Right panel: Strategic Actions ───────────────────────────────────────
    rp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, y_panel, right_w, panel_h,
    )
    rp.fill.solid()
    rp.fill.fore_color.rgb = CLR_LIGHT_BG
    rp.line.color.rgb = CLR_BORDER
    rp.line.width = Pt(0.5)

    # Action Badge
    needs_action = bool(exec_data.get("requires_action", False))
    act_text = "⚠ 需立即介入 (Action Required)" if needs_action else "✓ 常規監控 (Monitor)"
    act_clr = CLR_BAD if needs_action else CLR_GOOD
    _add_textbox(
        slide, right_left + Mm(5), y_panel + Mm(4), right_w - Mm(10), Mm(8),
        text=act_text, size_pt=10, bold=True, color=act_clr, alignment=PP_ALIGN.CENTER
    )

    # Check Directions
    _add_textbox(slide, right_left + Mm(5), y_panel + Mm(14), right_w - Mm(10), Mm(6), text="優先處理方向", size_pt=8, bold=True, color=CLR_SUBTITLE)
    check_dirs = list(exec_data.get("check_directions", []) or [])
    if not check_dirs:
        check_dirs = ["依據診斷結果執行現場複核。"]
    combined_dirs: List[str] = []
    if dashboard_layers:
        combined_dirs.extend(pptx_alarm_health_lines(dashboard_layers)[:4])
    combined_dirs.extend(check_dirs[:4])

    _add_list_box(
        slide, right_left + Mm(4), y_panel + Mm(22), right_w - Mm(8), panel_h - Mm(28),
        combined_dirs[:10], size_pt=7, color=CLR_BLACK
    )

    _add_footer(slide, timestamp_str, page_num, total_pages)


# ── P2: Core Diagnosis ───────────────────────────────────────────────────────

def _build_slide_core_diagnosis(
    prs: PptxPresentation,
    diagnostics: List[Dict[str, Any]],
    exec_data: Dict[str, Any],
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    *,
    dashboard_layers: Optional[Dict[str, Any]] = None,
):
    """
    P2 — 主要問題 / Core Diagnosis  (SMT SPI).
    Layout:
      Left: primary_feature (大字) + problem_type_label badge + primary anomaly
            summary + secondary warnings + min Cpk + stability verdict
      Right: risk sentence + recommended actions
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "主要問題 / Core Diagnosis", "主要特徵 · 問題型態 · 異常摘要 · 建議動作")

    content_top = MARGIN_TOP + Mm(28)
    content_h = Mm(130)
    left_w = Mm(133)
    gap = Mm(5)
    right_left = MARGIN_LEFT + left_w + gap
    right_w = CONTENT_WIDTH - left_w - gap

    # ── Left panel ──────────────────────────────────────────────────────────
    lp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT, content_top, left_w, content_h,
    )
    lp.fill.solid()
    lp.fill.fore_color.rgb = CLR_LIGHT_BG
    lp.line.color.rgb = CLR_BORDER
    lp.line.width = Pt(0.5)

    y = content_top + Mm(5)

    # 【主要特徵】large label
    primary_feature = str(exec_data.get("primary_feature", "—") or "—")
    problem_type_label = str(exec_data.get("problem_type_label", "") or "")
    _add_textbox(
        slide, MARGIN_LEFT + Mm(5), y,
        Mm(30), Mm(7),
        text="主要特徵", size_pt=8, bold=True, color=CLR_SUBTITLE,
    )
    _add_textbox(
        slide, MARGIN_LEFT + Mm(37), y,
        left_w - Mm(47), Mm(10),
        text=_ellipsize(primary_feature, 30), size_pt=13, bold=True, color=CLR_TITLE,
    )
    y += Mm(13)

    # 問題型態 badge (inline)
    has_pt = bool(problem_type_label and problem_type_label not in ("—", "Unknown", ""))
    if has_pt:
        pt_badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN_LEFT + Mm(5), y,
            left_w - Mm(10), Mm(10),
        )
        pt_badge.fill.solid()
        pt_badge.fill.fore_color.rgb = CLR_BADGE_INFO_BG
        pt_badge.line.color.rgb = CLR_BADGE_INFO_BORDER
        pt_badge.line.width = Pt(0.7)
        _add_textbox(
            slide, MARGIN_LEFT + Mm(7), y + Mm(2),
            Mm(26), Mm(6),
            text="問題型態", size_pt=7, bold=True, color=CLR_BADGE_INFO_TEXT,
        )
        _add_textbox(
            slide, MARGIN_LEFT + Mm(35), y + Mm(1),
            left_w - Mm(45), Mm(9),
            text=problem_type_label, size_pt=9, bold=True,
            color=CLR_ACCENT,
        )
        y += Mm(14)

    # 【主要異常】
    primary_diag = diagnostics[0] if diagnostics else {}
    severity_label_text = _severity_label(str(primary_diag.get("severity", "info")))
    sev_clr = _severity_color(str(primary_diag.get("severity", "info")))
    _add_textbox(
        slide, MARGIN_LEFT + Mm(5), y,
        left_w - Mm(10), Mm(7),
        text=f"【主要異常】{severity_label_text}",
        size_pt=9, bold=True, color=sev_clr,
    )
    y += Mm(8)
    primary_summary = _ellipsize(primary_diag.get("summary", "未偵測到異常診斷訊號"), 130)
    ps_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Mm(5), y,
        left_w - Mm(10), Mm(20),
    )
    ps_box.text_frame.word_wrap = True
    _p = ps_box.text_frame.paragraphs[0]
    _run = _p.add_run()
    _run.text = primary_summary
    _font(_run, 9, bold=False, color=CLR_BLACK)
    y += Mm(22)

    evidence_type = str(primary_diag.get("evidence_type", "") or "").strip()
    if evidence_type:
        _add_textbox(
            slide,
            MARGIN_LEFT + Mm(5),
            y,
            left_w - Mm(10),
            Mm(6),
            text=f"證據類型：{evidence_type}",
            size_pt=7,
            color=CLR_SUBTITLE,
        )
        y += Mm(7)

    # 【次要警示】
    secondary = [
        d for d in diagnostics[1:]
        if str(d.get("severity", "")) in ("warning", "error")
    ]
    if secondary:
        _add_textbox(
            slide, MARGIN_LEFT + Mm(5), y,
            left_w - Mm(10), Mm(7),
            text="【次要警示】", size_pt=9, bold=True, color=CLR_WARNING,
        )
        y += Mm(8)
        sec_lines = [
            _ellipsize(d.get("summary", "—"), 80)
            for d in secondary[:2]
        ]
        _add_list_box(
            slide, MARGIN_LEFT + Mm(5), y,
            left_w - Mm(10), Mm(18),
            sec_lines, size_pt=8, color=CLR_BLACK,
        )
        y += Mm(20)

    # 【最弱能力特徵】 + 【穩定性結論】on same row if space allows
    min_cpk_str = str(exec_data.get("min_cpk_str", "—"))
    stability = str(exec_data.get("stability_verdict", "受控"))
    stab_clr = (
        CLR_BAD if stability == "失控"
        else CLR_WARNING if stability in ("警示", "漂移", "偏移")
        else CLR_GOOD
    )
    # Box for min cpk
    cpk_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_LEFT + Mm(5), y, Mm(60), Mm(16))
    cpk_box.fill.solid()
    cpk_box.fill.fore_color.rgb = CLR_WHITE
    cpk_box.line.color.rgb = CLR_BORDER
    _add_textbox(slide, MARGIN_LEFT + Mm(7), y + Mm(1), Mm(56), Mm(5), text="最弱 Cpk", size_pt=7, bold=True, color=CLR_SUBTITLE)
    _add_textbox(slide, MARGIN_LEFT + Mm(7), y + Mm(6), Mm(56), Mm(8), text=min_cpk_str, size_pt=10, bold=True, color=_cpk_color(min_cpk_str))

    # Box for stability
    stab_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_LEFT + Mm(70), y, Mm(58), Mm(16))
    stab_box.fill.solid()
    stab_box.fill.fore_color.rgb = CLR_WHITE
    stab_box.line.color.rgb = CLR_BORDER
    _add_textbox(slide, MARGIN_LEFT + Mm(72), y + Mm(1), Mm(54), Mm(5), text="製程穩定性", size_pt=7, bold=True, color=CLR_SUBTITLE)
    _add_textbox(slide, MARGIN_LEFT + Mm(72), y + Mm(6), Mm(54), Mm(8), text=f"製程{stability}", size_pt=10, bold=True, color=stab_clr)

    # ── Right panel ──────────────────────────────────────────────────────────
    rp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, content_top, right_w, content_h,
    )
    rp.fill.solid()
    rp.fill.fore_color.rgb = CLR_LIGHT_BG
    rp.line.color.rgb = CLR_BORDER
    rp.line.width = Pt(0.5)

    ry = content_top + Mm(5)
    if dashboard_layers:
        l8_blob = "\n".join(pptx_layer8_diagnosis_lines(dashboard_layers))
        _add_textbox(
            slide,
            right_left + Mm(5),
            ry,
            right_w - Mm(10),
            Mm(38),
            text=l8_blob,
            size_pt=8,
            bold=False,
            color=CLR_BLACK,
        )
        ry += Mm(40)
    _add_textbox(
        slide, right_left + Mm(5), ry,
        right_w - Mm(10), Mm(7),
        text="【風險判讀】", size_pt=9, bold=True, color=CLR_ACCENT,
    )
    ry += Mm(8)
    risk_h = Mm(24) if dashboard_layers else Mm(32)
    risk_sent_box = slide.shapes.add_textbox(
        right_left + Mm(5), ry,
        right_w - Mm(10), risk_h,
    )
    risk_sent_box.text_frame.word_wrap = True
    _p = risk_sent_box.text_frame.paragraphs[0]
    _run = _p.add_run()
    _run.text = str(exec_data.get("risk_sentence", "—"))
    _font(_run, 9, bold=False, color=CLR_BLACK)
    ry = ry + risk_h + Mm(4)

    _add_textbox(
        slide, right_left + Mm(5), ry,
        right_w - Mm(10), Mm(7),
        text="【建議動作】", size_pt=9, bold=True, color=CLR_ACCENT,
    )
    ry += Mm(8)
    action_lines: List[str] = []
    seen_acts: set = set()
    for diag in diagnostics:
        for act in (diag.get("recommended_actions") or []):
            txt = str(act or "").strip()
            if txt and txt not in seen_acts:
                action_lines.append(txt)
                seen_acts.add(txt)
            if len(action_lines) >= 5:
                break
        if len(action_lines) >= 5:
            break
    if not action_lines:
        action_lines = ["依實際異常結果進行現場確認與追蹤。"]
    act_h = Mm(40) if dashboard_layers else Mm(52)
    _add_list_box(
        slide, right_left + Mm(5), ry,
        right_w - Mm(10), act_h,
        action_lines[:5], size_pt=8, color=CLR_BAD,
    )

    _add_footer(slide, timestamp_str, page_num, total_pages)


# ── P5: Multi-Signal Diagnosis ───────────────────────────────────────────────

def _build_slide_multi_signal_diagnosis(
    prs: PptxPresentation,
    msd: Dict[str, Any],
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    *,
    dashboard_layers: Optional[Dict[str, Any]] = None,
):
    """
    P5 — 多訊號關聯診斷 / Multi-Signal Diagnosis.
    Layout (A4 landscape):
      Left panel (~150 mm):
        • 異常訊號來源圖表 — table of signals (chart type / anomaly type / feature)
        • 異常之間的關聯  — correlation pattern label + engineering explanation
      Right panel (~108 mm):
        • 製程型態判定  — primary anomaly type badge
        • Chart types triggered list
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(
        slide,
        "多訊號關聯診斷 / Multi-Signal Diagnosis",
        "異常訊號來源 · 訊號關聯分析 · 製程型態判定",
    )

    content_top = MARGIN_TOP + Mm(28)
    content_h   = Mm(130)
    left_w      = Mm(150)
    gap         = Mm(5)
    right_left  = MARGIN_LEFT + left_w + gap
    right_w     = CONTENT_WIDTH - left_w - gap

    # ── Left panel ──────────────────────────────────────────────────────────
    lp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT, content_top, left_w, content_h,
    )
    lp.fill.solid()
    lp.fill.fore_color.rgb = CLR_LIGHT_BG
    lp.line.color.rgb = CLR_BORDER
    lp.line.width = Pt(0.5)

    y = content_top + Mm(5)
    if dashboard_layers:
        _add_textbox(
            slide,
            MARGIN_LEFT + Mm(5),
            y,
            left_w - Mm(10),
            Mm(26),
            text="\n".join(pptx_defect_structure_lines(dashboard_layers)),
            size_pt=7,
            bold=False,
            color=CLR_SUBTITLE,
        )
        y += Mm(28)

    # ── 異常訊號來源圖表 ────────────────────────────────────────────────────
    _add_textbox(
        slide, MARGIN_LEFT + Mm(5), y,
        left_w - Mm(10), Mm(7),
        text="【異常訊號來源圖表】", size_pt=9, bold=True, color=CLR_ACCENT,
    )
    y += Mm(8)

    signals: List[Dict[str, Any]] = msd.get("signals", [])
    sig_cap = 4 if dashboard_layers else 6
    if signals:
        # Header row
        col_x = [
            MARGIN_LEFT + Mm(5),
            MARGIN_LEFT + Mm(45),
            MARGIN_LEFT + Mm(90),
        ]
        col_w = [Mm(38), Mm(43), left_w - Mm(100)]
        hdr_labels = ["圖表類型", "異常型態", "特徵 Feature"]
        for hx, hw, hl in zip(col_x, col_w, hdr_labels):
            _add_textbox(
                slide, hx, y, hw, Mm(6),
                text=hl, size_pt=7, bold=True, color=CLR_SUBTITLE,
            )
        y += Mm(6)

        # Signal rows (max 4–6; tighter when儀表板區塊已佔高度)
        for sig in signals[:sig_cap]:
            sev = str(sig.get("severity", "warning"))
            row_clr = CLR_BAD if sev == "error" else CLR_WARNING
            vals = [
                sig.get("chart_type", "—"),
                sig.get("anomaly_type_zh", sig.get("anomaly_type", "—")),
                _ellipsize(sig.get("feature", "—"), 20),
            ]
            for hx, hw, val in zip(col_x, col_w, vals):
                _add_textbox(
                    slide, hx, y, hw, Mm(6),
                    text=str(val), size_pt=8, bold=False, color=row_clr,
                )
            y += Mm(6)
        if len(signals) > sig_cap:
            _add_textbox(
                slide, MARGIN_LEFT + Mm(5), y,
                left_w - Mm(10), Mm(5),
                text=f"…另有 {len(signals) - sig_cap} 個訊號（見異常診斷詳頁）",
                size_pt=7, bold=False, color=CLR_SUBTITLE,
            )
            y += Mm(6)
    else:
        _add_textbox(
            slide, MARGIN_LEFT + Mm(5), y,
            left_w - Mm(10), Mm(7),
            text="未偵測到 warning / error 等級的異常訊號。",
            size_pt=8, bold=False, color=CLR_SUBTITLE,
        )
        y += Mm(9)

    y += Mm(3)  # separator

    # ── 異常之間的關聯 ──────────────────────────────────────────────────────
    _add_textbox(
        slide, MARGIN_LEFT + Mm(5), y,
        left_w - Mm(10), Mm(7),
        text="【異常之間的關聯】", size_pt=9, bold=True, color=CLR_ACCENT,
    )
    y += Mm(8)

    correlation: Dict[str, Any] = msd.get("correlation", {})
    pattern_label  = str(correlation.get("pattern_label",  "—"))
    pattern_detail = str(correlation.get("pattern_detail", "—"))

    # Pattern label badge
    corr_badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT + Mm(5), y,
        left_w - Mm(10), Mm(8),
    )
    corr_badge.fill.solid()
    corr_badge.fill.fore_color.rgb = CLR_BADGE_WARNING_BG
    corr_badge.line.color.rgb      = CLR_BADGE_WARNING_BORDER
    corr_badge.line.width = Pt(0.7)
    _add_textbox(
        slide, MARGIN_LEFT + Mm(7), y + Mm(1),
        left_w - Mm(14), Mm(6),
        text=pattern_label, size_pt=9, bold=True, color=CLR_BADGE_WARNING_TEXT,
    )
    y += Mm(10)

    detail_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Mm(5), y,
        left_w - Mm(10), Mm(28),
    )
    detail_box.text_frame.word_wrap = True
    _p = detail_box.text_frame.paragraphs[0]
    _run = _p.add_run()
    _run.text = pattern_detail
    _font(_run, 8, bold=False, color=CLR_BLACK)

    # ── Right panel ──────────────────────────────────────────────────────────
    rp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, content_top, right_w, content_h,
    )
    rp.fill.solid()
    rp.fill.fore_color.rgb = CLR_LIGHT_BG
    rp.line.color.rgb = CLR_BORDER
    rp.line.width = Pt(0.5)

    ry = content_top + Mm(5)
    _add_textbox(slide, right_left + Mm(5), ry, right_w - Mm(10), Mm(7), text="【製程型態判定】", size_pt=9, bold=True, color=CLR_ACCENT)
    ry += Mm(8)

    primary_zh = str(msd.get("primary_anomaly_type_zh", "—"))
    # Prominent Verdict Badge
    _add_status_indicator(slide, right_left + Mm(4), ry, right_w - Mm(8), Mm(18), "主要診斷結果", primary_zh, CLR_ACCENT)
    ry += Mm(24)

    # KB Rules if exists
    kb_rules: List[Any] = msd.get("kb_matched_rules") or []
    if kb_rules:
        _add_textbox(slide, right_left + Mm(5), ry, right_w - Mm(10), Mm(6), text="知識庫匹配 (SPI-KB)", size_pt=7, bold=True, color=CLR_SUBTITLE)
        ry += Mm(6)
        kb_text = ""
        for r in kb_rules[:2]:
            if isinstance(r, dict):
                summary_txt = str(r.get("summary", "") or "")
                kb_text += f"• {r.get('rule_id')} ({r.get('match_confidence')}): {summary_txt[:40]}...\n"
        _add_textbox(slide, right_left + Mm(5), ry, right_w - Mm(10), Mm(15), text=kb_text, size_pt=7, color=CLR_BLACK)
        ry += Mm(18)

    matrix_raw = msd.get("diagnostic_evidence_matrix")
    matrix = matrix_raw if isinstance(matrix_raw, dict) else {}
    matrix_summary = matrix.get("summary") if isinstance(matrix.get("summary"), dict) else {}
    if matrix_summary:
        _add_textbox(slide, right_left + Mm(5), ry, right_w - Mm(10), Mm(6), text="組合證據矩陣", size_pt=8, bold=True, color=CLR_ACCENT)
        ry += Mm(7)
        readable_tabs = build_readable_diagnostic_tabs(matrix)
        readable_rows = [
            *readable_tabs.get("overview", [])[:1],
            *readable_tabs.get("chart_linkage", [])[:2],
        ]
        matrix_lines = [
            str(matrix_summary.get("coverage_line_zh") or "—"),
            f"矩陣結論: {_ellipsize(str(matrix_summary.get('verdict_zh') or '—'), 44)}",
        ]
        for item in readable_rows:
            if isinstance(item, dict):
                matrix_lines.append(
                    f"• {item.get('title', '—')}: "
                    f"{item.get('result_zh', '—')}；"
                    f"{_ellipsize(str(item.get('reason_zh') or '—'), 38)}"
                )
        _add_textbox(
            slide,
            right_left + Mm(5),
            ry,
            right_w - Mm(10),
            Mm(26),
            text="\n".join(matrix_lines),
            size_pt=7,
            color=CLR_BLACK,
        )
        ry += Mm(29)

    # Action Logic
    _add_textbox(slide, right_left + Mm(5), ry, right_w - Mm(10), Mm(6), text="處理建議方向", size_pt=8, bold=True, color=CLR_ACCENT)
    ry += Mm(7)
    _add_list_box(
        slide, right_left + Mm(5), ry, right_w - Mm(10), Mm(40),
        ["1. 與基準良率對比 (Verification)", "2. 跨特徵量測確認 (Cross-check)", "3. 依診斷調整製程參數"],
        size_pt=8, color=CLR_BLACK
    )

    _add_footer(slide, timestamp_str, page_num, total_pages)


# ── P6: Process Cause Hypothesis ─────────────────────────────────────────────

def _build_slide_cause_hypothesis(
    prs: PptxPresentation,
    msd: Dict[str, Any],
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    *,
    dashboard_layers: Optional[Dict[str, Any]] = None,
):
    """
    P6 — 製程原因推論 / Process Cause Hypothesis.
    Layout (A4 landscape):
      Left panel (~150 mm): cause hypotheses (category + description per row)
      Right panel (~108 mm): recommended check items (category + bullet items)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(
        slide,
        "製程原因推論 / Process Cause Hypothesis",
        "製程原因推論 · 建議檢查項目",
    )

    content_top = MARGIN_TOP + Mm(28)
    content_h   = Mm(130)
    left_w      = Mm(150)
    gap         = Mm(5)
    right_left  = MARGIN_LEFT + left_w + gap
    right_w     = CONTENT_WIDTH - left_w - gap

    # ── Left panel — 製程原因推論 ────────────────────────────────────────────
    lp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT, content_top, left_w, content_h,
    )
    lp.fill.solid()
    lp.fill.fore_color.rgb = CLR_LIGHT_BG
    lp.line.color.rgb = CLR_BORDER
    lp.line.width = Pt(0.5)

    y = content_top + Mm(5)
    _add_textbox(
        slide, MARGIN_LEFT + Mm(5), y,
        left_w - Mm(10), Mm(7),
        text="【可能原因推論 Process Cause Hypothesis】",
        size_pt=9, bold=True, color=CLR_ACCENT,
    )
    y += Mm(9)

    causes: List[Dict[str, str]] = msd.get("cause_hypotheses", [])
    if not causes:
        causes = [{"category": "—", "description": "未能從現有訊號形成可複核原因假設，請人工判讀。"}]

    # Category color palette (cycling)
    _CAT_COLORS = [
        CLR_BADGE_INFO_BG,
        CLR_BADGE_NEUTRAL_BG,
        CLR_BADGE_SUCCESS_BG,
        _rgb(PROCESS_ALARM_CARD_BG_WARNING),
        CLR_BADGE_ERROR_BG,
    ]
    _CAT_BORDER_COLORS = [
        CLR_ACCENT,
        _rgb(CHART_PALETTE_OFFSET_X_FILL),
        _rgb(CHART_PALETTE_HEIGHT_FILL),
        _rgb(CHART_PALETTE_AREA_FILL),
        _rgb(CHART_PALETTE_SOLDER_FILL),
    ]

    for i, cause in enumerate(causes[:5]):
        cat  = str(cause.get("category",    "—"))
        desc = str(cause.get("description", "—"))
        bg   = _CAT_COLORS[i % len(_CAT_COLORS)]
        bdr  = _CAT_BORDER_COLORS[i % len(_CAT_BORDER_COLORS)]

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN_LEFT + Mm(5), y,
            left_w - Mm(10), Mm(20),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg
        card.line.color.rgb      = bdr
        card.line.width = Pt(0.8)

        _add_textbox(
            slide, MARGIN_LEFT + Mm(7), y + Mm(1),
            left_w - Mm(14), Mm(7),
            text=cat, size_pt=9, bold=True, color=bdr,
        )
        desc_box = slide.shapes.add_textbox(
            MARGIN_LEFT + Mm(7), y + Mm(9),
            left_w - Mm(14), Mm(10),
        )
        desc_box.text_frame.word_wrap = True
        _p = desc_box.text_frame.paragraphs[0]
        _run = _p.add_run()
        _run.text = desc
        _font(_run, 8, bold=False, color=CLR_BLACK)

        y += Mm(22)

    # ── Right panel — 建議檢查項目 ───────────────────────────────────────────
    rp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, content_top, right_w, content_h,
    )
    rp.fill.solid()
    rp.fill.fore_color.rgb = CLR_LIGHT_BG
    rp.line.color.rgb = CLR_BORDER
    rp.line.width = Pt(0.5)

    ry = content_top + Mm(5)
    _add_textbox(
        slide, right_left + Mm(5), ry,
        right_w - Mm(10), Mm(7),
        text="【建議檢查項目 Check Items】",
        size_pt=9, bold=True, color=CLR_ACCENT,
    )
    ry += Mm(9)

    check_items: List[Dict[str, Any]] = msd.get("check_items", [])
    if not check_items:
        check_items = [{"category": "—", "items": ["請依核心診斷結果進行現場確認。"]}]
    if dashboard_layers:
        raw_l8 = dashboard_layers.get("layer_8_diagnosis")
        l8: Dict[str, Any] = raw_l8 if isinstance(raw_l8, dict) else {}
        act_l8 = str(l8.get("recommended_action_zh") or "").strip()
        if act_l8:
            check_items = [
                {"category": "儀表板建議對策", "items": [act_l8]},
                *check_items,
            ]

    for ci in check_items[:3]:
        cat   = str(ci.get("category", "—"))
        items = ci.get("items", [])

        _add_textbox(
            slide, right_left + Mm(5), ry,
            right_w - Mm(10), Mm(6),
            text=cat, size_pt=8, bold=True, color=CLR_ACCENT,
        )
        ry += Mm(7)

        for item in items[:3]:
            item_box = slide.shapes.add_textbox(
                right_left + Mm(7), ry,
                right_w - Mm(12), Mm(8),
            )
            item_box.text_frame.word_wrap = True
            _p = item_box.text_frame.paragraphs[0]
            _run = _p.add_run()
            _run.text = f"• {item}"
            _font(_run, 7, bold=False, color=CLR_BLACK)
            ry += Mm(8)

        ry += Mm(2)  # gap between categories

def _build_slide_process_capability_v2(
    prs: PptxPresentation,
    summary_data: dict,
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    *,
    dashboard_layers: Optional[Dict[str, Any]] = None,
):
    """
    P3a — 製程能力 / Process Capability.
    Redesigned to show feature comparison + spec summary.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "製程能力 / Process Capability", "Cpk · Ppk · 良率分析")

    content_top = MARGIN_TOP + Mm(26)
    per_measure = summary_data.get("per_measure", {}) if isinstance(summary_data, dict) else {}
    measures = ["Volume", "Area", "Height"]
    available_measures = [m for m in measures if m in per_measure]

    # 1. Spec & KPI Grid
    _add_textbox(slide, MARGIN_LEFT, content_top, Mm(50), Mm(6), "【各項能力指標明細】", size_pt=9, bold=True, color=CLR_ACCENT)
    
    n_cols = len(available_measures) + 1 if available_measures else 1
    table_shape = slide.shapes.add_table(5, n_cols, MARGIN_LEFT, content_top + Mm(8), CONTENT_WIDTH, Mm(40))
    table = table_shape.table
    _style_table_header(table, n_cols)
    _style_table_rows(table, 1, 5, n_cols)

    _set_cell_text(table.cell(0, 0), "指標Metric", size_pt=8, bold=True, alignment=PP_ALIGN.CENTER)
    from app.utils.constants import FEATURE_DISPLAY_NAMES
    for ci, m in enumerate(available_measures, 1):
        _set_cell_text(table.cell(0, ci), FEATURE_DISPLAY_NAMES.get(m, m), size_pt=8, bold=True, alignment=PP_ALIGN.CENTER)

    kpis = [("Cpk", "cap.cpk"), ("Ppk", "cap.ppk"), ("Yield %", "yield_pct"), ("PPM", "defect.ppm_total")]
    for ri, (label, path) in enumerate(kpis, 1):
        _set_cell_text(table.cell(ri, 0), label, size_pt=8, bold=True)
        for ci, m in enumerate(available_measures, 1):
            txt, _clr = _extract_metric_value(per_measure.get(m, {}), path)
            _set_cell_text(table.cell(ri, ci), txt, size_pt=8, alignment=PP_ALIGN.CENTER)

    # 2. Logic Note & Warnings
    y_note = content_top + Mm(55)
    _add_textbox(slide, MARGIN_LEFT, y_note, CONTENT_WIDTH, Mm(6), "【能力評鑑說明】", size_pt=9, bold=True, color=CLR_ACCENT)
    notes = [
        "• 能力充足 (Cpk > 1.33)：製程寬量足夠，容許輕微波動。",
        "• 能力邊緣 (Cpk 1.0~1.33)：需加強穩定性監控，防止不合格產出。",
        "• 能力不足 (Cpk < 1.00)：規格寬度不足 or 製程分布過寬，建議檢討 Spec 或改善製程。",
        "• Cpk 與 Ppk 顯著差異 (>10%)：表示存在顯著的長期變異 (Long-term drift)。"
    ]
    if dashboard_layers:
        notes.extend(pptx_kpi_capability_lines(dashboard_layers)[:2])

    _add_list_box(slide, MARGIN_LEFT + Mm(4), y_note + Mm(7), CONTENT_WIDTH - Mm(8), Mm(40), notes, size_pt=8, color=CLR_BLACK)

    _add_footer(slide, timestamp_str, page_num, total_pages)


def _build_slide_process_stability(
    prs: PptxPresentation,
    diagnostics: List[Dict[str, Any]],
    exec_data: Dict[str, Any],
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    *,
    dashboard_layers: Optional[Dict[str, Any]] = None,
):
    """
    P3b — 製程穩定性 / Process Stability.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "製程穩定性 / Process Stability", "SPC · CUSUM · EWMA · 趨勢分析")

    content_top = MARGIN_TOP + Mm(26)
    stability = str(exec_data.get("stability_verdict", "受控"))
    stab_clr = CLR_BAD if stability == "失控" else CLR_WARNING if stability in ("警示", "偏移", "漂移") else CLR_GOOD

    # 1. Stability Status Badge
    _add_status_indicator(slide, MARGIN_LEFT, content_top, CONTENT_WIDTH, Mm(20), "穩定性判定 (Stability Verdict)", f"製程{stability}", stab_clr)
    
    # 2. Key Stability Insights
    y_panel = content_top + Mm(25)
    _add_textbox(slide, MARGIN_LEFT, y_panel, Mm(60), Mm(6), "【異常診斷摘要】", size_pt=9, bold=True, color=CLR_ACCENT)
    
    lines = []
    if dashboard_layers:
        lines.extend(pptx_stability_dashboard_lines(dashboard_layers)[:3])
    
    spc_diag = _find_diagnostic_by_keywords(diagnostics, ["cusum", "ewma", "run", "趨勢", "管制", "shift", "drift"])
    if spc_diag:
        lines.append(f"• 核心現象：{_ellipsize(spc_diag.get('summary', ''), 120)}")
        lines.extend([f"  - {ev}" for ev in (spc_diag.get("evidence_lines") or [])[:3]])
    else:
        lines.append("• 本批次未偵測到顯著的統計學漂移或偏移訊號。")
    
    _add_list_box(slide, MARGIN_LEFT + Mm(4), y_panel + Mm(8), CONTENT_WIDTH - Mm(8), Mm(45), lines[:8], size_pt=8, color=CLR_BLACK)

    # 3. Stability Guidance
    _add_textbox(slide, MARGIN_LEFT, y_panel + Mm(55), Mm(60), Mm(6), "【控制策略建議】", size_pt=9, bold=True, color=CLR_SUBTITLE)
    guidance = [
        "• 失控/警示：建議檢查印刷機參數、刮刀、鋼網與錫膏狀態。",
        "• 趨勢/漂移：可能為錫膏乾涸或環境溫度變化，需確認刮刀壓力。",
        "• 受控：建議維持現行管制圖寬度，常規監控即可。"
    ]
    _add_list_box(slide, MARGIN_LEFT + Mm(4), y_panel + Mm(62), CONTENT_WIDTH - Mm(8), Mm(25), guidance, size_pt=8, color=CLR_BLACK)

    _add_footer(slide, timestamp_str, page_num, total_pages)


# ── P3c: Process Risk (enhanced) ─────────────────────────────────────────────

def _build_slide_process_risk_v2(
    prs: PptxPresentation,
    exec_data: Dict[str, Any],
    risk_assessment: Dict[str, Any],
    process: Dict[str, Any],
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    *,
    report_context: Optional[Dict[str, Any]] = None,
):
    """
    P3c — 製程風險 / Process Risk (enhanced).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "製程風險 / Process Risk", "能力 + 穩定性 + 分布 → 綜合風險判定")

    content_top = MARGIN_TOP + Mm(26)
    ra = risk_assessment if isinstance(risk_assessment, dict) else {}
    risk_level = str(exec_data.get("risk_level", "LOW")).upper()
    risk_display = str(exec_data.get("risk_display", risk_level))
    risk_clr = CLR_BAD if risk_level == "HIGH" else CLR_WARNING if risk_level == "MEDIUM" else CLR_GOOD

    # 1. Overall Risk Status
    _add_status_indicator(slide, MARGIN_LEFT, content_top, CONTENT_WIDTH, Mm(24), "整體製程風險等級 (Overall Risk)", risk_display, risk_clr)
    _add_textbox(
        slide,
        MARGIN_LEFT + Mm(6),
        content_top + Mm(18),
        CONTENT_WIDTH - Mm(12),
        Mm(6),
        text=f"整體風險：{risk_display}",
        size_pt=8,
        bold=True,
        color=risk_clr,
    )
    
    # 2. Key Synthesis
    ry = content_top + Mm(30)
    _add_textbox(slide, MARGIN_LEFT, ry, Mm(60), Mm(6), "【數據綜合判斷】", size_pt=9, bold=True, color=CLR_ACCENT)
    ry += Mm(8)
    
    risk_sentence = str(exec_data.get("risk_sentence", "—"))
    s_box = slide.shapes.add_textbox(MARGIN_LEFT + Mm(4), ry, CONTENT_WIDTH - Mm(8), Mm(30))
    s_box.text_frame.word_wrap = True
    _p = s_box.text_frame.paragraphs[0]
    _run = _p.add_run()
    _run.text = risk_sentence
    _font(_run, 11, bold=True, color=CLR_TITLE)
    
    # 3. Anomaly Summary Table
    ry += Mm(32)
    _add_textbox(slide, MARGIN_LEFT, ry, Mm(60), Mm(6), "【異常規模統計】", size_pt=9, bold=True, color=CLR_SUBTITLE)
    ry += Mm(7)
    
    error_count = int(ra.get("error_count", exec_data.get("error_count", 0)) or 0)
    warning_count = int(ra.get("warning_count", exec_data.get("warning_count", 0)) or 0)
    high_priority = int(ra.get("high_priority_count", exec_data.get("high_priority_count", 0)) or 0)
    
    n_c = 3
    tbl = slide.shapes.add_table(2, n_c, MARGIN_LEFT + Mm(4), ry, CONTENT_WIDTH - Mm(8), Mm(16)).table
    _style_table_header(tbl, n_c)
    _set_cell_text(tbl.cell(0,0), "Error (嚴重)", size_pt=8, alignment=PP_ALIGN.CENTER)
    _set_cell_text(tbl.cell(0,1), "Warning (警示)", size_pt=8, alignment=PP_ALIGN.CENTER)
    _set_cell_text(tbl.cell(0,2), "High Priority (優先)", size_pt=8, alignment=PP_ALIGN.CENTER)
    _set_cell_text(tbl.cell(1,0), str(error_count), size_pt=10, bold=True, color=CLR_BAD, alignment=PP_ALIGN.CENTER)
    _set_cell_text(tbl.cell(1,1), str(warning_count), size_pt=10, bold=True, color=CLR_WARNING, alignment=PP_ALIGN.CENTER)
    _set_cell_text(tbl.cell(1,2), str(high_priority), size_pt=10, bold=True, color=CLR_ACCENT, alignment=PP_ALIGN.CENTER)
    _add_textbox(
        slide,
        MARGIN_LEFT + Mm(6),
        ry + Mm(17),
        CONTENT_WIDTH - Mm(12),
        Mm(6),
        text=f"高優先訊號：{high_priority}",
        size_pt=8,
        color=CLR_SUBTITLE,
    )
    trust_lines = _trust_status_lines(report_context or {})
    if not trust_lines:
        trust_lines = [
            "統計/能力：可信：資料直接計算",
            "根因/對策：需複核：規則推論",
        ]
    _add_textbox(
        slide,
        MARGIN_LEFT,
        ry + Mm(27),
        Mm(60),
        Mm(6),
        text="【報告可信度狀態】",
        size_pt=9,
        bold=True,
        color=CLR_SUBTITLE,
    )
    _add_list_box(
        slide,
        MARGIN_LEFT + Mm(4),
        ry + Mm(35),
        CONTENT_WIDTH - Mm(8),
        Mm(24),
        trust_lines[:4],
        size_pt=8,
        color=CLR_BLACK,
    )

    _add_footer(slide, timestamp_str, page_num, total_pages)


# ── P9: Background Info (merged workorder + spec) ────────────────────────────

def _build_slide_background_info(
    prs: PptxPresentation,
    wo_master: dict,
    wo_spec: dict,
    report_context: dict,
    timestamp_str: str,
    page_num: int,
    total_pages: int,
):
    """
    P9 — Background Info (合併工單 + 規格).
    Replaces old slides 1+2 with a compact two-column layout.
    Moved to near-end so it doesn't delay the diagnosis section.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(
        slide,
        "背景資訊 / Background Info",
        "工單資訊 · 產品規格",
    )
    # Backward-compatible section anchors kept for legacy report consumers/tests.
    _add_textbox(
        slide,
        MARGIN_LEFT,
        MARGIN_TOP + Mm(23),
        CONTENT_WIDTH,
        Mm(4),
        text="1. Product & Work Order Information | 2. Control Specification",
        size_pt=6,
        color=CLR_SUBTITLE,
    )

    content_top = MARGIN_TOP + Mm(28)
    content_h = Mm(130)
    left_w = Mm(126)
    gap = Mm(5)
    right_left = MARGIN_LEFT + left_w + gap
    right_w = CONTENT_WIDTH - left_w - gap

    # ── Left: Work Order (compact) ──────────────────────────────────────────
    lp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT, content_top, left_w, content_h,
    )
    lp.fill.solid()
    lp.fill.fore_color.rgb = CLR_LIGHT_BG
    lp.line.color.rgb = CLR_BORDER
    lp.line.width = Pt(0.5)

    _add_textbox(
        slide, MARGIN_LEFT + Mm(5), content_top + Mm(4),
        left_w - Mm(10), Mm(7),
        text="工單資訊 (Work Order)", size_pt=10, bold=True, color=CLR_ACCENT,
    )

    inferred = report_context.get("inferred_context", {}) if isinstance(report_context, dict) else {}
    inferred_line = inferred.get("line_name") if isinstance(inferred, dict) else None
    product_part_no = _resolve_workorder_field(
        wo_master, report_context, key="product_part_no", fallback_key="product_part_no",
    )
    line_name = wo_master.get("line_name")
    if _is_missing_value(line_name) and not _is_missing_value(inferred_line):
        line_name = inferred_line

    wo_fields = [
        ("供應商製令工單", _report_value(wo_master.get("supplier_work_order_no", "—"))),
        ("醫電製令工單", _report_value(wo_master.get("outsource_work_order_no") or wo_master.get("work_order_no", "—"))),
        ("產品名稱", _report_value(wo_master.get("product_name", "—"))),
        ("產品料號", _report_value(product_part_no)),
        ("供應商", _report_value(wo_master.get("supplier", "—"))),
        ("批量", _report_value(wo_master.get("batch_qty", "—"))),
        ("線別", _report_value(line_name)),
        ("生產日期", _report_value(wo_master.get("production_date", "—"))),
        ("錫膏批號", _report_value(wo_master.get("solder_paste_lot", "—"))),
        ("PCB DateCode", _report_value(wo_master.get("pcb_datecode", "—"))),
    ]
    wy = content_top + Mm(14)
    for lbl, val in wo_fields:
        _add_textbox(
            slide, MARGIN_LEFT + Mm(6), wy, Mm(36), Mm(7),
            text=f"{lbl}：", size_pt=9, bold=True, color=CLR_SUBTITLE,
        )
        _add_textbox(
            slide, MARGIN_LEFT + Mm(44), wy, left_w - Mm(50), Mm(7),
            text=_report_value(val), size_pt=9, color=CLR_BLACK,
        )
        wy += Mm(8)

    # Completeness note
    required_keys_map = {
        "supplier_work_order_no": "供應商製令工單", "outsource_work_order_no": "醫電製令工單",
        "product_name": "產品名稱",
        "product_part_no": "產品料號", "supplier": "供應商", "batch_qty": "批量",
    }
    missing = [
        label for key, label in required_keys_map.items()
        if _is_missing_value(wo_master.get(key))
    ]
    if missing:
        _add_textbox(
            slide, MARGIN_LEFT + Mm(6), wy + Mm(2), left_w - Mm(12), Mm(8),
            text=f"主檔完整度：必填欄位缺失 {len(missing)}/5。缺失：{', '.join(missing)}",
            size_pt=7, color=CLR_WARNING,
        )
    entry = report_context.get("coordinate_registry_entry", {}) if isinstance(report_context, dict) else {}
    if (
        isinstance(entry, dict)
        and not _is_missing_value(entry.get("product_part_no"))
        and _is_missing_value(wo_master.get("product_part_no"))
    ):
        _add_textbox(
            slide,
            MARGIN_LEFT + Mm(6),
            wy + Mm(10),
            left_w - Mm(12),
            Mm(7),
            text="產品料號已由座標註冊表自動回填",
            size_pt=7,
            color=CLR_SUBTITLE,
        )

    # ── Right: Control Spec (compact) ───────────────────────────────────────
    rp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, content_top, right_w, content_h,
    )
    rp.fill.solid()
    rp.fill.fore_color.rgb = CLR_LIGHT_BG
    rp.line.color.rgb = CLR_BORDER
    rp.line.width = Pt(0.5)

    _add_textbox(
        slide, right_left + Mm(5), content_top + Mm(4),
        right_w - Mm(10), Mm(7),
        text="量測特徵規格 (LSL / USL / Target)", size_pt=10, bold=True, color=CLR_ACCENT,
    )

    spec_rows = []
    for key, label in [("volume", "Volume [%]"), ("area", "Area [%]"), ("height", "Height [%]")]:
        s = wo_spec.get(key, {})
        if s:
            spec_rows.append((
                label,
                _report_value(s.get("usl", "—")),
                _report_value(s.get("lsl", "—")),
                _report_value(s.get("target", "—")),
            ))
    if not spec_rows:
        spec_rows.append(("（未設定規格）", "—", "—", "—"))

    n_rows = len(spec_rows) + 1
    tbl_w = right_w - Mm(10)
    tbl = slide.shapes.add_table(
        n_rows, 4, right_left + Mm(5), content_top + Mm(14), tbl_w, Mm(8 * n_rows),
    ).table
    col_widths_spec = [Mm(34), Mm(16), Mm(16), Mm(16)]
    for ci, cw in enumerate(col_widths_spec[:4]):
        if ci < len(tbl.columns):
            tbl.columns[ci].width = cw
    for ci, hdr in enumerate(["量測特徵", "USL", "LSL", "Target"]):
        _set_cell_text(tbl.cell(0, ci), hdr, size_pt=8, bold=True, alignment=PP_ALIGN.CENTER)
    _style_table_header(tbl, 4)
    for ri, (feat, usl, lsl, tgt) in enumerate(spec_rows, 1):
        _set_cell_text(tbl.cell(ri, 0), feat, size_pt=8, bold=True)
        _set_cell_text(tbl.cell(ri, 1), usl, size_pt=8, alignment=PP_ALIGN.CENTER)
        _set_cell_text(tbl.cell(ri, 2), lsl, size_pt=8, alignment=PP_ALIGN.CENTER)
        _set_cell_text(tbl.cell(ri, 3), tgt, size_pt=8, alignment=PP_ALIGN.CENTER)

    # Spec source note
    profile = report_context.get("product_spec_profile", {}) if isinstance(report_context, dict) else {}
    if isinstance(profile, dict) and profile:
        spec_note = (
            f"規格來源：{profile.get('product_name', '—')} / "
            f"Stencil: {profile.get('stencil_type', '—')}"
        )
        _add_textbox(
            slide, right_left + Mm(5), content_top + Mm(14) + Mm(8 * n_rows) + Mm(3),
            right_w - Mm(10), Mm(8),
            text=spec_note, size_pt=7, color=CLR_SUBTITLE,
        )
        _add_textbox(
            slide, right_left + Mm(5), content_top + Mm(14) + Mm(8 * n_rows) + Mm(10),
            right_w - Mm(10), Mm(8),
            text=(
                f"Main Thickness: {_format_optional_mm(profile.get('thickness_main'))}, "
                f"Precision Thickness: {_format_optional_mm(profile.get('thickness_precision'))}"
            ),
            size_pt=7,
            color=CLR_SUBTITLE,
        )

    _add_footer(slide, timestamp_str, page_num, total_pages)


def _build_slide_statistics(
    prs: PptxPresentation,
    summary_data: dict,
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    *,
    dashboard_layers: Optional[Dict[str, Any]] = None,
):
    """Section 3: Statistics summary KPIs + per-feature metrics table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "3. Statistics Summary", "製程統計分析摘要")

    per_measure = summary_data.get("per_measure", {})
    process = summary_data.get("process", {})

    content_top = MARGIN_TOP + Mm(28)

    # ── KPI bar ─────────────────────────────────────────────────────────────
    kpi_data = [
        ("樣本數 (N)", _get_sample_n(per_measure), CLR_BLACK),
        ("整體良率 (%)", _fmt_pct(process.get("overall_yield_pct")),
         _yield_color(process.get("overall_yield_pct"))),
        ("最弱 Cpk", _fmt_cpk(process.get("min_cpk"), process.get("min_cpk_measure")),
         _cpk_color(process.get("min_cpk"))),
        ("製程判定", process.get("verdict", "—"),
         _verdict_color(process.get("verdict", "—"))),
    ]

    kpi_card_w = Mm(62)
    kpi_card_h = Mm(28)
    kpi_x = int(MARGIN_LEFT)
    for label, value, clr in kpi_data:
        # KPI card background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            kpi_x, content_top, kpi_card_w, kpi_card_h,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = CLR_LIGHT_BG
        bg.line.color.rgb = CLR_BORDER
        bg.line.width = Pt(0.5)

        # Value
        _add_textbox(
            slide, kpi_x + Mm(4), content_top + Mm(3), kpi_card_w - Mm(8), Mm(14),
            text=str(value), size_pt=9, bold=True, color=clr, alignment=PP_ALIGN.CENTER,
        )
        # Label
        _add_textbox(
            slide, kpi_x + Mm(4), content_top + Mm(17), kpi_card_w - Mm(8), Mm(8),
            text=label, size_pt=9, color=CLR_SUBTITLE, alignment=PP_ALIGN.CENTER,
        )
        kpi_x += kpi_card_w + Mm(4)

    defect_combined = process.get("defect_combined", {}) if isinstance(process, dict) else {}
    dpmo_event = _to_float(defect_combined.get("dpmo_combined_event")) if isinstance(defect_combined, dict) else None
    if dpmo_event is not None:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            content_top + kpi_card_h + Mm(2),
            CONTENT_WIDTH,
            Mm(6),
            text=f"Combined DPMO(Event): {dpmo_event:.1f}",
            size_pt=8,
            color=CLR_SUBTITLE,
        )

    if dashboard_layers:
        eng_lines = "\n".join(pptx_engineering_data_lines(dashboard_layers))
        _add_textbox(
            slide, MARGIN_LEFT, content_top + kpi_card_h + Mm(2),
            CONTENT_WIDTH, Mm(14),
            text=eng_lines, size_pt=7, color=CLR_SUBTITLE,
        )
        table_top = content_top + kpi_card_h + Mm(18)
    else:
        table_top = content_top + kpi_card_h + Mm(10)
    measures = ["Volume", "Area", "Height"]
    available_measures = [m for m in measures if m in per_measure]

    metric_rows_def = [
        ("Cp", "cap.cp"),
        ("Cpk", "cap.cpk"),
        ("Pp", "cap.pp"),
        ("Ppk", "cap.ppk"),
        ("Mean (平均值)", "dist.mean"),
        ("Stdev (標準差)", "cap.sigma_lt"),
        ("N (樣本數)", "n"),
        ("Yield (良率 %)", "yield_pct"),
        ("PPM (Total)", "defect.ppm_total"),
        ("DPMO (Feature)", "defect.dpmo_feature"),
        ("Zbench ST", "defect.zbench_st"),
        ("Zbench LT", "defect.zbench_lt"),
        ("Cpk CI", "defect.cpk_ci"),
        ("Cpk CI Method", "defect.cpk_ci_method"),
    ]

    n_cols = 1 + len(available_measures)
    n_rows = 1 + len(metric_rows_def)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        MARGIN_LEFT, table_top,
        CONTENT_WIDTH, Mm(7 * n_rows),
    )
    table = table_shape.table

    # Column widths
    first_col_w = Mm(50)
    remaining = CONTENT_WIDTH - first_col_w
    data_col_w = int(remaining // len(available_measures)) if available_measures else int(remaining)
    table.columns[0].width = first_col_w
    for ci in range(1, n_cols):
        table.columns[ci].width = data_col_w

    # Header row
    from app.utils.constants import FEATURE_DISPLAY_NAMES
    _set_cell_text(table.cell(0, 0), "統計指標", size_pt=9, bold=True, alignment=PP_ALIGN.CENTER)
    for ci, m in enumerate(available_measures, 1):
        display = FEATURE_DISPLAY_NAMES.get(m, m)
        _set_cell_text(table.cell(0, ci), display, size_pt=9, bold=True, alignment=PP_ALIGN.CENTER)
    _style_table_header(table, n_cols)

    # Data rows
    _style_table_rows(table, 1, n_rows, n_cols)
    for ri, (label, path) in enumerate(metric_rows_def, 1):
        _set_cell_text(table.cell(ri, 0), label, size_pt=8, bold=True)
        for ci, m in enumerate(available_measures, 1):
            text, _clr = _extract_metric_value(per_measure.get(m, {}), path)
            _set_cell_text(table.cell(ri, ci), text, size_pt=8, alignment=PP_ALIGN.CENTER)

    _add_footer(slide, timestamp_str, page_num, total_pages)


def _build_slide_diagnostic(
    prs: PptxPresentation,
    diagnostic: Dict[str, Any],
    slide_index: int,
    timestamp_str: str,
    page_num: int,
    total_pages: int,
    *,
    title_text: Optional[str] = None,
    subtitle_text: str = "Anomaly Diagnosis",
):
    """Build anomaly diagnostic slide with associated evidence/chart context."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    severity = str(diagnostic.get("severity", "info") or "info")
    sev_label = _severity_label(severity)
    sev_color = _severity_color(severity)
    feature_label = _ellipsize(diagnostic.get("feature_label", "") or "", 28)
    summary = _ellipsize(diagnostic.get("summary", "") or "—", 180)
    chart_bytes = diagnostic.get("chart_bytes")
    chart_title = _ellipsize(diagnostic.get("chart_title", "") or "相關圖表", 64)
    chart_missing_reason = _ellipsize(
        diagnostic.get("chart_missing_reason", "") or "此異常目前無可用圖表輸出。",
        72,
    )
    observable_charts = [_ellipsize(item, 32) for item in (diagnostic.get("observable_charts", []) or [])[:2]]
    evidence_lines = [_ellipsize(item, 42) for item in (diagnostic.get("evidence_lines", []) or [])[:3]]
    ipc_lines = [_ellipsize(item, 52) for item in (diagnostic.get("ipc_lines", []) or [])[:2]]
    evidence_type = _ellipsize(diagnostic.get("evidence_type", "") or "統計計算 / 規則推論", 48)
    action_lines = [
        _ellipsize(item, 40)
        for item in ((diagnostic.get("recommended_actions", []) or ["依本頁證據與圖表進行現場複核。"])[:2])
    ]

    if title_text:
        _add_slide_title(slide, title_text, subtitle_text)
    else:
        _add_slide_title(slide, f"異常診斷詳細分析 / Anomaly Detail Analysis ({slide_index})", "單項特徵現象與建議")

    content_top = MARGIN_TOP + Mm(28)
    left_w = Mm(120)
    right_left = MARGIN_LEFT + left_w + Mm(6)
    right_w = CONTENT_WIDTH - left_w - Mm(6)
    content_h = Mm(130)

    # ── Left panel: Root cause analysis ─────────────────────────────────────
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT, content_top, left_w, content_h,
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = CLR_LIGHT_BG
    left_panel.line.color.rgb = CLR_BORDER
    left_panel.line.width = Pt(0.5)

    # Severity badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT + Mm(4), content_top + Mm(4), Mm(32), Mm(9),
    )
    badge.fill.solid()
    if severity == "error":
        badge.fill.fore_color.rgb = CLR_BADGE_ERROR_BG
    elif severity == "warning":
        badge.fill.fore_color.rgb = CLR_BADGE_WARNING_BG
    else:
        badge.fill.fore_color.rgb = CLR_BADGE_NEUTRAL_BG
    badge.line.fill.background()

    _add_textbox(
        slide, MARGIN_LEFT + Mm(6), content_top + Mm(5), Mm(28), Mm(6),
        text=sev_label, size_pt=10, bold=True, color=sev_color,
    )
    if feature_label:
        _add_textbox(
            slide, MARGIN_LEFT + Mm(40), content_top + Mm(5), left_w - Mm(46), Mm(6),
            text=f"量測特徵：{feature_label}", size_pt=8, bold=True, color=CLR_SUBTITLE,
        )

    y_pos = content_top + Mm(18)
    _add_textbox(
        slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(8),
        text="異常摘要", size_pt=9, bold=True, color=CLR_ACCENT,
    )
    y_pos += Mm(8)
    summary_box = _add_textbox(
        slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(24),
        text=summary, size_pt=9, color=CLR_BLACK,
    )
    summary_box.text_frame.word_wrap = True
    y_pos += Mm(26)

    if observable_charts:
        _add_textbox(
            slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(7),
            text=f"對應圖表：{', '.join(observable_charts[:2])}", size_pt=8, color=CLR_SUBTITLE,
        )
        y_pos += Mm(7)

    _add_textbox(
        slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(7),
        text=f"證據類型：{evidence_type}", size_pt=8, color=CLR_SUBTITLE,
    )
    y_pos += Mm(7)

    _add_textbox(
        slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(7),
        text="建議檢查/動作", size_pt=9, bold=True, color=CLR_ACCENT,
    )
    y_pos += Mm(8)
    _add_list_box(
        slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(16),
        action_lines[:2], size_pt=8, color=CLR_BAD,
    )
    y_pos += Mm(18)

    if evidence_lines:
        _add_textbox(
            slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(7),
            text="關鍵證據", size_pt=9, bold=True, color=CLR_ACCENT,
        )
        y_pos += Mm(8)
        _add_colored_list_box(
            slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(18),
            evidence_lines[:4], size_pt=8, default_color=CLR_BLACK, color_resolver=_evidence_line_color,
        )
        y_pos += Mm(20)

    if ipc_lines:
        _add_textbox(
            slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(7),
            text="IPC 參考", size_pt=9, bold=True, color=CLR_ACCENT,
        )
        y_pos += Mm(8)
        _add_list_box(
            slide, MARGIN_LEFT + Mm(6), y_pos, left_w - Mm(12), Mm(14),
            ipc_lines[:2], size_pt=7, color=CLR_SUBTITLE,
        )

    # ── Right panel: Chart image ────────────────────────────────────────────
    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, content_top, right_w, content_h,
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = CLR_WHITE
    right_panel.line.color.rgb = CLR_BORDER
    right_panel.line.width = Pt(0.5)

    _add_textbox(
        slide, right_left + Mm(5), content_top + Mm(4), right_w - Mm(10), Mm(8),
        text=chart_title, size_pt=10, bold=True, color=CLR_TITLE,
    )

    if chart_bytes:
        _add_picture_contain(
            slide,
            chart_bytes,
            right_left + Mm(4),
            content_top + Mm(15),
            right_w - Mm(8),
            content_h - Mm(19),
        )
    else:
        # No chart placeholder
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            right_left + Mm(4), content_top + Mm(15), right_w - Mm(8), content_h - Mm(19),
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = CLR_LIGHT_BG
        placeholder.line.color.rgb = CLR_BORDER
        _add_textbox(
            slide, right_left + Mm(18), content_top + Mm(58), right_w - Mm(36), Mm(18),
            text="此異常無可用圖表輸出", size_pt=12, color=CLR_SUBTITLE, alignment=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide, right_left + Mm(16), content_top + Mm(74), right_w - Mm(32), Mm(16),
            text=chart_missing_reason, size_pt=8, color=CLR_SUBTITLE, alignment=PP_ALIGN.CENTER,
        )

    _add_footer(slide, timestamp_str, page_num, total_pages)


# ═══════════════════════════════════════════════════════════════════════════════
# Value extraction helpers
# ═══════════════════════════════════════════════════════════════════════════════

def _get_sample_n(per_measure: dict) -> str:
    """Get max N across all measures."""
    max_n = 0
    for m_data in per_measure.values():
        n_raw = m_data.get("n")
        try:
            n_val = coerce_int(n_raw)
        except (TypeError, ValueError):
            continue
        if n_val > max_n:
            max_n = n_val
    return f"{max_n:,}" if max_n > 0 else "—"


def _fmt_pct(value) -> str:
    parsed = _to_float(value)
    if parsed is None:
        return "—"
    normalized = parsed * 100.0 if 0.0 <= parsed <= 1.0 else parsed
    return f"{normalized:.2f}%"


def _fmt_cpk(cpk, measure=None) -> str:
    if cpk is None:
        return "—"
    try:
        base = f"{float(cpk):.3f}"
        if measure:
            return f"{base} ({measure})"
        return base
    except (ValueError, TypeError):
        return "—"


def _to_float(value: Any) -> Optional[float]:
    """Parse numeric-like values safely using shared utility."""
    return safe_float(value)


def _extract_metric_value(pm: dict, key: str) -> Tuple[str, RGBColor]:
    """Extract a metric value and its color from per_measure data."""
    if not pm:
        return "—", CLR_BLACK

    cap = (pm.get("cap") or {}).get("statistics") or {}
    dist = (pm.get("dist") or {}).get("statistics") or {}
    defect = pm.get("defect") or {}

    val = None
    fmt = "—"
    color = CLR_BLACK

    if key == "cap.cp":
        val = _to_float(cap.get("cp"))
        if val is not None:
            fmt = f"{val:.3f}"
            color = _cpk_color(val)
    elif key == "cap.cpk":
        val = _to_float(cap.get("cpk"))
        if val is not None:
            fmt = f"{val:.3f}"
            color = _cpk_color(val)
    elif key == "cap.pp":
        val = _to_float(cap.get("pp"))
        if val is not None:
            fmt = f"{val:.3f}"
            color = _cpk_color(val)
    elif key == "cap.ppk":
        val = _to_float(cap.get("ppk"))
        if val is not None:
            fmt = f"{val:.3f}"
            color = _cpk_color(val)
    elif key == "dist.mean":
        val = _to_float(dist.get("mean"))
        if val is not None:
            fmt = f"{val:.2f}"
    elif key == "cap.sigma_lt":
        val = _to_float(cap.get("sigma_lt"))
        if val is None:
            val = _to_float(dist.get("std"))
        if val is not None:
            fmt = f"{val:.3f}"
    elif key == "n":
        val = _to_float(pm.get("n"))
        if val is not None:
            fmt = f"{int(val):,}"
    elif key == "yield_pct":
        val = _to_float(pm.get("yield_pct"))
        if val is not None:
            yield_pct = val * 100.0 if 0.0 <= val <= 1.0 else val
            fmt = f"{yield_pct:.2f}%"
            color = _yield_color(yield_pct)
    elif key == "defect.ppm_total":
        val = _to_float(defect.get("ppm_total"))
        if val is not None:
            fmt = f"{val:.1f}"
            if val <= 100:
                color = CLR_GOOD
            elif val <= 10000:
                color = CLR_WARNING
            else:
                color = CLR_BAD
    elif key == "defect.dpmo_feature":
        val = _to_float(defect.get("dpmo_feature"))
        if val is not None:
            fmt = f"{val:.1f}"
            if val <= 100:
                color = CLR_GOOD
            elif val <= 10000:
                color = CLR_WARNING
            else:
                color = CLR_BAD
    elif key == "defect.zbench_st":
        val = _to_float(defect.get("zbench_st"))
        if val is not None:
            fmt = f"{val:.3f}"
            color = _cpk_color(val / 3.0)
    elif key == "defect.zbench_lt":
        val = _to_float(defect.get("zbench_lt"))
        if val is not None:
            fmt = f"{val:.3f}"
            color = _cpk_color(val / 3.0)
    elif key == "defect.cpk_ci":
        ci_text = str(defect.get("cpk_ci", "") or "").strip()
        if ci_text:
            fmt = ci_text
    elif key == "defect.cpk_ci_method":
        ci_method = str(defect.get("cpk_ci_method", "") or "").strip()
        if ci_method:
            if ci_method == "N/A":
                fmt = ci_method
            elif "Bissell" in ci_method:
                fmt = "Bissell 95%"
            else:
                fmt = ci_method

    return fmt, color


# ═══════════════════════════════════════════════════════════════════════════════
# Main build function
# ═══════════════════════════════════════════════════════════════════════════════

def build_pptx_report(
    wo_master: Dict[str, Any],
    wo_spec: Dict[str, Any],
    summary_data: Dict[str, Any],
    diagnostics: List[Dict[str, Any]],
    output_path: str,
    analysis_payload: Optional[Dict[str, Any]] = None,
    report_context: Optional[Dict[str, Any]] = None,
    template_type: str = "engineering",
    chart_ids_to_export: Optional[List[str]] = None,
    progress_callback: Optional[Callable[[int, str], None]] = None,
    render_chart_fn: Optional[Callable[..., Optional[bytes]]] = None,
) -> Tuple[bool, Optional[str]]:
    """
    Build a 12-section-framework PPTX report and save to output_path.

    Args:
        wo_master: Work order master data dict.
        wo_spec: Workorder spec dict (volume/area/height with usl/lsl/target).
        summary_data: Output from compute_summary() — per_measure, relation, process.
        diagnostics: Diagnostic entry list with summary, actions, chart title and image bytes.
        analysis_payload: Full analysis payload for richer section evidence.
        report_context: Export context (relation_meta, filter context, registry/profile info).
        output_path: File path to save the PPTX.
        template_type: Sole supported value is ``engineering`` (legacy values ignored for slide copy).

    Returns:
        (True, None) on success, (False, error_message) on failure.
    """
    try:
        if progress_callback:
            progress_callback(5, "初始化 PPTX 引擎...")
        logger.debug("build_pptx_report template_type=%s", template_type)
        timestamp_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not diagnostics:
            diagnostics = [{
                "summary": "目前未觸發異常診斷規則；本批資料未發現需要額外開立異常頁的高風險訊號。",
                "severity": "info",
                "chart_title": "異常診斷狀態",
                "chart_bytes": None,
                "observable_charts": [],
                "evidence_lines": [
                    "狀態：未觸發異常規則 (anomaly rule)",
                    "輸出策略：核心 12 章節 + 圖表證據擴充頁",
                ],
                "ipc_lines": [],
                "recommended_actions": [
                    "若需更細部證據，請於匯出頁勾選相關圖表後重新匯出 PPTX。",
                    "持續觀察 Cpk、Yield 與趨勢圖是否出現新漂移。",
                ],
                "evidence_type": "統計計算 / 規則推論",
            }]
        analysis_payload = analysis_payload or {}
        report_context = report_context or {}
        selected_chart_ids = _dedupe_chart_ids(chart_ids_to_export or [])
        selected_features = [
            str(feature).strip()
            for feature in (report_context.get("selected_features", []) if isinstance(report_context, dict) else [])
            if str(feature).strip()
        ]
        available_features = [
            str(feature).strip()
            for feature in (report_context.get("available_features", []) if isinstance(report_context, dict) else [])
            if str(feature).strip()
        ]
        if not selected_features:
            selected_features = available_features[:1]
        coverage_by_id = _coverage_item_by_id(report_context) if isinstance(report_context, dict) else {}
        if progress_callback:
            progress_callback(10, "正在渲染圖表證據 (這可能需要一點時間)...")
        chart_evidence_items = _render_chart_evidence_items(
            selected_chart_ids=selected_chart_ids,
            analysis_payload=analysis_payload if isinstance(analysis_payload, dict) else {},
            selected_features=selected_features,
            available_features=available_features,
            coverage_by_id=coverage_by_id,
            render_chart_fn=render_chart_fn,
        )
        chart_gallery_pages = (len(chart_evidence_items) + 3) // 4
        coverage_pages = _coverage_pages_count(report_context) if isinstance(report_context, dict) else 0
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        process = summary_data.get("process", {}) if isinstance(summary_data, dict) else {}
        dashboard_layers = extract_dashboard_layers(process)
        has_dashboard = bool(dashboard_layers)
        dashboard_layers_for_slides: Optional[Dict[str, Any]] = (
            dashboard_layers if has_dashboard else None
        )
        has_process_layers_slide = bool(
            isinstance(report_context, dict) and report_context.get("process_diagnosis_report")
        )
        # Base layout = 13 slides, plus optional process-layer bridge slide.
        base_pages = 13 + (1 if has_process_layers_slide else 0)
        total_pages = base_pages + len(diagnostics) + chart_gallery_pages + coverage_pages
        relation = summary_data.get("relation", {}) if isinstance(summary_data, dict) else {}
        per_measure = summary_data.get("per_measure", {}) if isinstance(summary_data, dict) else {}

        # Pre-compute multi-signal diagnosis for P5/P6.
        msd = _msd.run_multi_signal_diagnosis(diagnostics)
        matrix = (
            report_context.get("diagnostic_evidence_matrix", {})
            if isinstance(report_context, dict)
            else {}
        )
        if isinstance(matrix, dict) and matrix:
            matrix_relation_raw = matrix.get("relation")
            matrix_relation: Dict[str, Any] = (
                matrix_relation_raw if isinstance(matrix_relation_raw, dict) else {}
            )
            matrix_signals = matrix_relation.get("signals") if isinstance(matrix_relation.get("signals"), list) else []
            matrix_corr = matrix_relation.get("correlation") if isinstance(matrix_relation.get("correlation"), dict) else {}
            matrix_causes = matrix_relation.get("cause_hypotheses") if isinstance(matrix_relation.get("cause_hypotheses"), list) else []
            matrix_checks = matrix_relation.get("check_items") if isinstance(matrix_relation.get("check_items"), list) else []
            msd["diagnostic_evidence_matrix"] = matrix
            if matrix_signals:
                msd["signals"] = matrix_signals
                distinct_charts = {str(s.get("chart_type") or "") for s in matrix_signals if isinstance(s, dict)}
                msd["has_multi_signal"] = len(distinct_charts) >= 2
            if matrix_corr:
                msd["correlation"] = matrix_corr
            if matrix_causes:
                msd["cause_hypotheses"] = matrix_causes
            if matrix_checks:
                msd["check_items"] = matrix_checks

        # Pre-compute executive summary data for P1/P2/P3b/P3c.
        risk_assessment = (
            report_context.get("risk_assessment", {})
            if isinstance(report_context, dict)
            else {}
        )
        if not (isinstance(risk_assessment, dict) and risk_assessment):
            risk_assessment = report_risk.build_risk_assessment(
                process=process if isinstance(process, dict) else {},
                diagnostics=diagnostics,
            )
        if not isinstance(risk_assessment, dict):
            risk_assessment = {}
        high_priority_count = sum(
            1 for item in diagnostics
            if isinstance(item, dict) and str(item.get("priority", "")).strip().lower() == "high"
        )
        if high_priority_count > 0:
            risk_assessment["high_priority_count"] = high_priority_count
        verdict_text = str(process.get("verdict", "") or "").strip()
        if high_priority_count > 0 or verdict_text in {"不可接受", "待改善"}:
            risk_assessment["level"] = "HIGH"
            risk_assessment["level_display"] = "高風險 (High)"
        exec_data = _build_exec_summary_data(
            summary_data or {},
            diagnostics,
            risk_assessment,
            report_context=report_context if isinstance(report_context, dict) else {},
        )

        # Running page counter (slides are built in diagnosis-first order)
        _pn = 0

        def _next_pn() -> int:
            nonlocal _pn
            _pn += 1
            # Report progress roughly based on page count vs total
            if progress_callback:
                prog = 20 + int((_pn / (total_pages + 1)) * 75)
                progress_callback(prog, f"產生第 {_pn}/{total_pages} 頁報告...")
            return _pn

        # ── P1: Executive Summary ────────────────────────────────────────────
        _build_slide_executive_summary(
            prs,
            exec_data,
            timestamp_str,
            _next_pn(),
            total_pages,
            dashboard_layers=dashboard_layers_for_slides,
            report_context=report_context if isinstance(report_context, dict) else None,
        )

        # ── P2: Core Diagnosis ───────────────────────────────────────────────
        _build_slide_core_diagnosis(
            prs,
            diagnostics,
            exec_data,
            timestamp_str,
            _next_pn(),
            total_pages,
            dashboard_layers=dashboard_layers_for_slides,
        )

        if has_process_layers_slide:
            pdr_raw = report_context.get("process_diagnosis_report") if isinstance(report_context, dict) else {}
            pdr = pdr_raw if isinstance(pdr_raw, dict) else {}
            _build_slide_bullet_section(
                prs,
                title_text="製程診斷架構（四層）",
                subtitle_text="Decision / Diagnosis / Evidence / Data",
                lines=_format_process_diagnosis_report_lines(pdr),
                timestamp_str=timestamp_str,
                page_num=_next_pn(),
                total_pages=total_pages,
                size_pt=9,
            )

        # ── P3a: Process Capability ──────────────────────────────────────────
        _build_slide_process_capability_v2(
            prs,
            summary_data,
            timestamp_str,
            _next_pn(),
            total_pages,
            dashboard_layers=dashboard_layers_for_slides,
        )

        # ── P3b: Process Stability ───────────────────────────────────────────
        _build_slide_process_stability(
            prs,
            diagnostics,
            exec_data,
            timestamp_str,
            _next_pn(),
            total_pages,
            dashboard_layers=dashboard_layers_for_slides,
        )

        # ── P5: Multi-Signal Diagnosis ───────────────────────────────────────
        _build_slide_multi_signal_diagnosis(
            prs,
            msd,
            timestamp_str,
            _next_pn(),
            total_pages,
            dashboard_layers=dashboard_layers_for_slides,
        )

        # ── P6: Process Cause Hypothesis ─────────────────────────────────────
        _build_slide_cause_hypothesis(
            prs,
            msd,
            timestamp_str,
            _next_pn(),
            total_pages,
            dashboard_layers=dashboard_layers_for_slides,
        )

        # ── P3c: Process Risk ────────────────────────────────────────────────
        _build_slide_process_risk_v2(
            prs,
            exec_data,
            risk_assessment,
            process,
            timestamp_str,
            _next_pn(),
            total_pages,
            report_context=report_context if isinstance(report_context, dict) else None,
        )

        # ── P4: Chart Evidence Gallery (moved before detail pages) ───────────
        if chart_gallery_pages > 0:
            for gallery_idx in range(chart_gallery_pages):
                start = gallery_idx * 4
                end = start + 4
                _build_slide_chart_evidence_gallery(
                    prs,
                    chart_items=chart_evidence_items[start:end],
                    gallery_index=gallery_idx + 1,
                    gallery_total=chart_gallery_pages,
                    timestamp_str=timestamp_str,
                    page_num=_next_pn(),
                    total_pages=total_pages,
                )

        if coverage_pages > 0:
            for coverage_idx in range(coverage_pages):
                _build_slide_chart_evidence_coverage(
                    prs,
                    report_context=report_context,
                    page_index=coverage_idx + 1,
                    page_total=coverage_pages,
                    timestamp_str=timestamp_str,
                    page_num=_next_pn(),
                    total_pages=total_pages,
                )

        # ── P5: Anomaly Diagnosis detail pages ───────────────────────────────
        for idx, diagnostic in enumerate(diagnostics, start=1):
            _build_slide_diagnostic(
                prs,
                diagnostic,
                slide_index=idx,
                timestamp_str=timestamp_str,
                page_num=_next_pn(),
                total_pages=total_pages,
                title_text=f"異常診斷詳頁 ({idx}/{len(diagnostics)}) — Anomaly Diagnosis",
                subtitle_text="異常診斷與建議",
            )

        # ── P6: Distribution Analysis (evidence zone) ────────────────────────
        dist_diag = _find_diagnostic_by_keywords(
            diagnostics, ["normal", "常態", "shapiro", "distribution", "分布"]
        )
        dist_lines = []
        if dist_diag:
            dist_lines.append(f"分布診斷：{_ellipsize(dist_diag.get('summary', '—'), 96)}")
            dist_lines.extend([str(line) for line in (dist_diag.get("evidence_lines") or [])[:4]])
        else:
            dist_lines.append("目前無顯著分布異常告警；請仍確認常態假設適用性。")
        dist_lines.append("若非常態，建議改用非常態能力分析或分層評估。")
        _build_slide_bullet_section(
            prs,
            title_text="分布分析 / Distribution Analysis",
            subtitle_text="常態性 · 分布形態",
            lines=dist_lines[:7],
            timestamp_str=timestamp_str,
            page_num=_next_pn(),
            total_pages=total_pages,
            colored=True,
        )

        # ── P7: Spatial Analysis ─────────────────────────────────────────────
        spatial_diag = _find_diagnostic_by_keywords(
            diagnostics, ["spatial", "空間", "edge", "pcb", "pad"]
        )
        relation_meta = (
            report_context.get("relation_meta", {}) if isinstance(report_context, dict) else {}
        )
        spatial_payload = (
            analysis_payload.get("spatial", {}) if isinstance(analysis_payload, dict) else {}
        )
        spatial_stats = (
            spatial_payload.get("statistics", {}) if isinstance(spatial_payload, dict) else {}
        )
        spatial_meta = (
            spatial_payload.get("metadata", {}) if isinstance(spatial_payload, dict) else {}
        )
        has_coordinate_scope = _has_coordinate_scope(
            report_context if isinstance(report_context, dict) else {},
            spatial_payload if isinstance(spatial_payload, dict) else {},
        )
        if not has_coordinate_scope:
            spatial_lines = [
                "本批資料未提供座標欄位，空間分析未納入判讀。",
                "章節狀態：未納入：資料缺失",
                "證據類型：未納入",
                "不輸出任何座標匹配率、空間點數或空間熱圖作為有效證據。",
            ]
        else:
            spatial_lines = ["目標：定位 PCB / Pad / Component 的偏差群聚與熱區。"]
            if dashboard_layers_for_slides:
                raw_l4 = dashboard_layers_for_slides.get("layer_4_defect_structure")
                l4s: Dict[str, Any] = raw_l4 if isinstance(raw_l4, dict) else {}
                spatial_lines.append(
                    f"儀表板 Top5 異常位號：{top_refdes_line(l4s.get('top_oos_refdes'))}"
                )
            match_rate = _to_float((relation_meta or {}).get("match_rate"))
            if match_rate is not None:
                spatial_lines.append(f"座標關聯成功率：{match_rate:.1f}%")
            match_count = relation_meta.get("match_count")
            unmatch_count = relation_meta.get("unmatch_count")
            if match_count is not None or unmatch_count is not None:
                spatial_lines.append(
                    f"匹配統計：match={match_count if match_count is not None else '—'}, "
                    f"unmatch={unmatch_count if unmatch_count is not None else '—'}"
                )
            points = _to_float((spatial_stats or {}).get("points"))
            if points is not None:
                spatial_lines.append(f"空間有效點數：{int(points):,}")
            if isinstance(spatial_meta, dict) and spatial_meta:
                spatial_mode = str(spatial_meta.get("mode", "")).strip()
                if spatial_mode:
                    spatial_lines.append(f"空間分析模式：{spatial_mode}")
            if spatial_diag:
                spatial_lines.append(f"空間異常：{_ellipsize(spatial_diag.get('summary', '—'), 96)}")
                missing_reason = str(spatial_diag.get("chart_missing_reason", "")).strip()
                if missing_reason:
                    spatial_lines.append(f"空間圖限制：{missing_reason}")
            else:
                spatial_lines.append("目前無空間異常提示；建議定期輸出熱圖做區域比對。")
            spatial_lines.append("輸出建議：PCB 熱圖 + Pad 分布 + Component 排名。")
        _build_slide_bullet_section(
            prs,
            title_text="空間分析 / Spatial Analysis (PCB / Pad / Component)",
            subtitle_text="熱區 · 群聚 · 邊緣異常",
            lines=spatial_lines[:7],
            timestamp_str=timestamp_str,
            page_num=_next_pn(),
            total_pages=total_pages,
        )

        # ── P8: Variation Source Analysis ────────────────────────────────────
        action_lines: List[str] = []
        seen_actions: set = set()
        for diag in diagnostics:
            for action in (diag.get("recommended_actions") or []):
                text = str(action).strip()
                if not text or text in seen_actions:
                    continue
                action_lines.append(text)
                seen_actions.add(text)
                if len(action_lines) >= 5:
                    break
            if len(action_lines) >= 5:
                break
        variation_lines = [
            "變異來源拆解框架：4M1E (Machine / Material / Method / Man / Environment)",
            "優先確認設備參數、原料批次、環境條件與治工具狀態。",
        ]
        if action_lines:
            variation_lines.append("建議改善動作（來自異常診斷）：")
            variation_lines.extend(action_lines[:4])
        else:
            variation_lines.append("目前缺少可歸因動作建議，請補強異常診斷資料。")
        corr_vol_area = _to_float((relation or {}).get("corr_vol_area"))
        corr_vol_height = _to_float((relation or {}).get("corr_vol_height"))
        corr_area_height = _to_float((relation or {}).get("corr_area_height"))
        if corr_vol_area is not None or corr_vol_height is not None or corr_area_height is not None:
            variation_lines.append(
                "關聯係數："
                f"Vol-Area={f'{corr_vol_area:.3f}' if corr_vol_area is not None else '—'}, "
                f"Vol-Height={f'{corr_vol_height:.3f}' if corr_vol_height is not None else '—'}, "
                f"Area-Height={f'{corr_area_height:.3f}' if corr_area_height is not None else '—'}"
            )
        pareto_data = (
            ((analysis_payload.get("pareto") or {}).get("components") or [])
            if isinstance(analysis_payload, dict) else []
        )
        if isinstance(pareto_data, list) and pareto_data:
            top_components = []
            for comp in pareto_data[:3]:
                if not isinstance(comp, dict):
                    continue
                cid = str(comp.get("component_id", "")).strip()
                rate = _to_float(comp.get("abnormal_rate"))
                if not cid:
                    continue
                top_components.append(f"{cid}({rate:.1%})" if rate is not None else cid)
            if top_components:
                variation_lines.append(f"TOP 異常元件：{', '.join(top_components)}")
        _build_slide_bullet_section(
            prs,
            title_text="變異來源分析 / Variation Source Analysis",
            subtitle_text="4M1E · 相關性 · 改善動作",
            lines=variation_lines[:7],
            timestamp_str=timestamp_str,
            page_num=_next_pn(),
            total_pages=total_pages,
        )

        # ── P9: Background Info (merged workorder + spec, moved to near-end) ──
        _build_slide_background_info(
            prs, wo_master or {}, wo_spec or {}, report_context,
            timestamp_str, _next_pn(), total_pages,
        )

        # ── P10: Statistics Summary (slide title; appendix zone — not a desktop UI page) ─────
        _build_slide_statistics(
            prs,
            summary_data or {},
            timestamp_str,
            _next_pn(),
            total_pages,
            dashboard_layers=dashboard_layers_for_slides,
        )

        appendix_lines = [
            "A. Data Source：量測資料、工單、規格、分析 payload",
            "B. 指標定義：Cp/Cpk/Pp/Ppk、Yield、PPM、OOC",
            "C. 檢定方法：常態性、趨勢與異常規則",
            "D. 圖表清單：見圖表證據覆蓋表（狀態含已輸出/不相容/無資料/本次排除/渲染失敗）",
            "E. 追蹤項目：VERIFY / UNKNOWN 欄位補完",
        ]
        available_features = [name for name in ("Volume", "Area", "Height") if name in per_measure]
        if available_features:
            appendix_lines.append(f"F. 本次分析特徵：{', '.join(available_features)}")
        filter_lines = _collect_filter_context_lines(report_context)
        if filter_lines:
            appendix_lines.append(f"G. 範圍：{'; '.join(filter_lines)}")
        scope_lines = _data_scope_lines(report_context) if isinstance(report_context, dict) else []
        if scope_lines:
            appendix_lines.extend(scope_lines[:3])
        if isinstance(relation_meta, dict) and relation_meta:
            duplicate_count = len((relation_meta.get("duplicate_coord_refdes") or []))
            unmatched_sample = relation_meta.get("unmatched_refdes_sample") or []
            appendix_lines.append(
                f"H. 關聯品質：duplicate_refdes={duplicate_count}, unmatched_sample_n={len(unmatched_sample)}"
            )
        product_part_no = _resolve_workorder_field(
            wo_master or {},
            report_context,
            key="product_part_no",
            fallback_key="product_part_no",
        )
        required_fields = {
            "供應商製令工單": wo_master.get("supplier_work_order_no"),
            "醫電製令工單": wo_master.get("outsource_work_order_no") or wo_master.get("work_order_no"),
            "產品名稱": wo_master.get("product_name"),
            "產品料號": product_part_no,
            "供應商": wo_master.get("supplier"),
            "批量": wo_master.get("batch_qty"),
        }
        missing_required = [label for label, value in required_fields.items() if _is_missing_value(value)]
        if missing_required:
            appendix_lines.append(f"I. 主檔缺失：{', '.join(missing_required)}（請於工單主檔補值）")
        # ── P11: Appendix ────────────────────────────────────────────────────
        _build_slide_bullet_section(
            prs,
            title_text="附錄 / Appendix",
            subtitle_text="定義 · 方法 · 追蹤清單",
            lines=appendix_lines[:8],
            timestamp_str=timestamp_str,
            page_num=_next_pn(),
            total_pages=total_pages,
        )

        if progress_callback:
            progress_callback(98, "儲存報告檔案...")
        prs.save(output_path)
        if progress_callback:
            progress_callback(100, "報告匯出完成")
        return True, None
    except (AttributeError, KeyError, TypeError, ValueError, RuntimeError, OSError) as e:
        logger.exception("PPTX 報告生成失敗: %s", e)
        return False, str(e)
