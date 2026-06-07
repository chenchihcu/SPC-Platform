"""PPTX export uses same filtered DF/spec as golden; stats table matches compute_summary."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.analytics.summary_engine import compute_summary
from app.data.session_store import SessionStore
from app.services.report_service import ReportService
from tests.release_validation.helpers.golden_scenario import load_joined_normal_baseline


def _all_slide_text(prs: Presentation) -> str:
    """Flatten visible text; include table cell text (python-pptx table shapes often have empty shape.text)."""
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                tbl = shape.table
                for ri in range(len(tbl.rows)):
                    for ci in range(len(tbl.columns)):
                        cell_txt = tbl.cell(ri, ci).text.strip()
                        if cell_txt:
                            parts.append(cell_txt)
            elif hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def test_pptx_statistics_summary_matches_golden_compute_summary(monkeypatch, tmp_path: Path, golden_root: Path) -> None:
    """SessionStore singleton populated from normal_baseline join; deck contains summary mean/yield for Volume."""
    store = SessionStore()
    store.clear()
    try:
        _sdir, _manifest, joined_df, spec = load_joined_normal_baseline(golden_root)
        store.joined_df = joined_df.copy()
        store.meas_meta = {"is_valid": True, "missing_required": []}
        store.coord_meta = {"is_valid": True, "missing_required": []}
        store.relation_meta = {"match_rate": 100.0, "unmatch_count": 0}
        store.workorder_spec = spec
        store.selected_features = ["Volume"]
        store.workorder_master = {
            "work_order_no": "RV-PHASE5",
            "product_name": "ReleaseValidation",
            "batch_qty": str(len(joined_df)),
        }

        ref = compute_summary(joined_df, spec)
        vol = ref["per_measure"]["Volume"]
        dist_mean = float((vol.get("dist") or {}).get("statistics", {}).get("mean"))
        mean_snippet = f"{dist_mean:.2f}"

        # Contract: per_measure yield_pct is 0-100 (same as summary_engine._yield_pct).
        display_yield = float(vol["yield_pct"])
        yield_snippet = f"{display_yield:.2f}%"

        monkeypatch.setattr(
            "app.services.chart_render.render_chart_to_png_bytes",
            lambda *args, **kwargs: None,
        )
        out = tmp_path / "rv-golden-report.pptx"
        ok, err = ReportService().generate_pptx_report(str(out))
        assert ok, err
        assert out.is_file()

        deck = _all_slide_text(Presentation(str(out)))
        assert mean_snippet in deck
        assert yield_snippet in deck
    finally:
        store.clear()
