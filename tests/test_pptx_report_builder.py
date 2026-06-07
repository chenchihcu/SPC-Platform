import base64

from pptx import Presentation

from app.services import pptx_report_builder as builder
from app.services.report_context import build_chart_evidence_coverage
from app.ui.theme.tokens import (
    ACCENT_PRIMARY,
    BG_BLOCK,
    RPT_SURFACE,
    TEXT_PRIMARY,
    TEXT_SECONDARY,
)


def _deck_text(prs: Presentation) -> str:
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            parts.append(cell.text)
            elif hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def test_pptx_palette_uses_shared_ui_theme_tokens() -> None:
    assert builder.CLR_ACCENT == builder._rgb(ACCENT_PRIMARY)
    assert builder.CLR_LIGHT_BG == builder._rgb(RPT_SURFACE)
    assert builder.CLR_WHITE == builder._rgb(BG_BLOCK)
    assert builder.CLR_TITLE == builder._rgb(TEXT_PRIMARY)
    assert builder.CLR_SUBTITLE == builder._rgb(TEXT_SECONDARY)


def test_evidence_line_color_cpk_thresholds() -> None:
    assert builder._evidence_line_color("Cpk: 1.80") == builder.CLR_GOOD
    assert builder._evidence_line_color("Cpk: 1.50") == builder.CLR_WARNING
    assert builder._evidence_line_color("Cpk: 1.10") == builder.CLR_BAD


def test_evidence_line_color_highlights_common_risk_metrics() -> None:
    assert builder._evidence_line_color("Yield: 99.50%") == builder.CLR_GOOD
    assert builder._evidence_line_color("Yield: 96.00%") == builder.CLR_WARNING
    assert builder._evidence_line_color("PPM (Total): 50000") == builder.CLR_BAD
    assert builder._evidence_line_color("DPMO：5000") == builder.CLR_WARNING
    assert builder._evidence_line_color("OOC Ratio: 12.5%") == builder.CLR_BAD
    assert builder._evidence_line_color("OOS 比率 (規格)：0.08%") == builder.CLR_BAD
    assert builder._evidence_line_color("Variance Ratio: 1.50x") == builder.CLR_WARNING
    assert builder._evidence_line_color("p-value: 0.01") == builder.CLR_BAD
    assert builder._evidence_line_color("Is Normal: False") == builder.CLR_BAD
    assert builder._evidence_line_color("Is Normal: True") == builder.CLR_GOOD


def test_evidence_line_color_defaults_to_black_for_neutral_lines() -> None:
    assert builder._evidence_line_color("Sample N: 120") == builder.CLR_BLACK


def test_process_diagnosis_report_lines_use_readable_evidence() -> None:
    lines = builder._format_process_diagnosis_report_lines(
        {
            "A_decision": {
                "process_verdict": "待改善",
                "core_diagnosis_zh": "目前證據顯示需要補強判讀。",
                "risk_level": "HIGH",
            },
            "B_diagnosis": {
                "scope": "全批",
                "distribution_shape": "Normal",
                "process_patterns": ["漂移"],
                "hypothesis_domain": "process",
            },
            "C_evidence": {
                "bridge_zh": "以白話證據列呈現主要判讀。",
                "combination_coverage_zh": "組合覆蓋 1/1 (100%)",
                "readable_evidence": [
                    {
                        "title": "中心偏移",
                        "result_zh": "不支持此假設",
                        "reason_zh": "目前未看到穩定的均值偏離證據。",
                        "evidence_zh": "Xbar-R / Volume",
                        "next_action_zh": "先確認漂移或資料完整性。",
                        "source_zh": "evidence_matrix.cells",
                    }
                ],
                "top_evidence": [
                    {
                        "chart_name": "raw-chart-name",
                        "feature_set": ["Volume"],
                        "metric_snapshot": "raw metric only",
                    }
                ],
            },
            "D_data": {"bridge_zh": "資料範圍完整。"},
        }
    )

    joined = "\n".join(lines)
    assert "不支持此假設" in joined
    assert "目前未看到穩定的均值偏離證據" in joined
    assert "反證" not in joined
    assert "raw-chart-name / Volume / raw metric only" not in joined


def test_build_pptx_report_applies_colored_evidence_runs(tmp_path) -> None:
    output_path = tmp_path / "builder-colored-evidence.pptx"
    ok, err = builder.build_pptx_report(
        wo_master={},
        wo_spec={},
        summary_data={},
        diagnostics=[
            {
                "summary": "Cpk 低於門檻",
                "severity": "warning",
                "chart_title": "分布能力",
                "chart_bytes": None,
                "observable_charts": [],
                "evidence_lines": ["Cpk: 0.82", "Sample N: 120"],
                "ipc_lines": [],
                "recommended_actions": ["確認設備參數穩定性"],
            }
        ],
        output_path=str(output_path),
    )

    assert ok, err
    prs = Presentation(str(output_path))
    diag_slide = None
    for slide in prs.slides:
        blob = "\n".join(
            s.text for s in slide.shapes if hasattr(s, "text") and s.text
        )
        if "Cpk: 0.82" in blob and "Sample N: 120" in blob:
            diag_slide = slide
            break
    assert diag_slide is not None

    cpk_color = None
    sample_n_color = None
    for shape in diag_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text or ""
                if "Cpk: 0.82" in text:
                    cpk_color = run.font.color.rgb
                if "Sample N: 120" in text:
                    sample_n_color = run.font.color.rgb

    assert cpk_color == builder.CLR_BAD
    assert sample_n_color == builder.CLR_BLACK


def test_extract_metric_value_formats_yield_with_percent_and_string_inputs() -> None:
    text, color = builder._extract_metric_value(
        {
            "yield_pct": "99.34",
            "cap": {"statistics": {"cpk": "1.20"}},
            "dist": {"statistics": {"mean": "100.1"}},
            "defect": {"ppm_total": "88.2"},
        },
        "yield_pct",
    )
    assert text == "99.34%"
    assert color == builder.CLR_GOOD

    ratio_text, ratio_color = builder._extract_metric_value(
        {"yield_pct": 0.9934},
        "yield_pct",
    )
    assert ratio_text == "99.34%"
    assert ratio_color == builder.CLR_GOOD


def test_fmt_pct_accepts_ratio_and_percent_values() -> None:
    assert builder._fmt_pct(99.34) == "99.34%"
    assert builder._fmt_pct(0.9934) == "99.34%"


def test_yield_color_accepts_ratio_and_percent_values() -> None:
    assert builder._yield_color(99.34) == builder.CLR_GOOD
    assert builder._yield_color(0.9934) == builder.CLR_GOOD
    assert builder._yield_color(0.96) == builder.CLR_WARNING
    assert builder._yield_color(0.90) == builder.CLR_BAD


def test_get_sample_n_accepts_numeric_like_strings() -> None:
    assert (
        builder._get_sample_n(
            {
                "Volume": {"n": "47660.0"},
                "Area": {"n": "47659"},
                "Height": {"n": None},
            }
        )
        == "47,660"
    )


def test_extract_metric_value_formats_n_from_string_values() -> None:
    text, color = builder._extract_metric_value({"n": "47660.0"}, "n")
    assert text == "47,660"
    assert color == builder.CLR_BLACK


def test_extract_metric_value_formats_cpk_ci_method_compactly() -> None:
    text, color = builder._extract_metric_value(
        {
            "defect": {
                "cpk_ci_method": "Bissell approximation (NIST/AIAG convention, two-sided 95%)"
            }
        },
        "defect.cpk_ci_method",
    )
    assert text == "Bissell 95%"
    assert color == builder.CLR_BLACK


def test_build_pptx_report_includes_contextual_completeness_data(tmp_path) -> None:
    output_path = tmp_path / "builder-contextual-completeness.pptx"
    ok, err = builder.build_pptx_report(
        wo_master={
            "work_order_no": "WO-CTX-01",
            "product_name": "X3000",
            "batch_qty": "20",
            "product_part_no": "",
        },
        wo_spec={
            "volume": {"usl": "150", "lsl": "70", "target": "100"},
            "area": {"usl": "150", "lsl": "70", "target": "100"},
            "height": {"usl": "140", "lsl": "70", "target": "100"},
        },
        summary_data={
            "process": {
                "overall_yield_pct": 99.1,
                "min_cpk": 1.32,
                "min_cpk_measure": "Height",
                "verdict": "待改善",
                "defect_combined": {
                    "dpmo_combined_event": 1200.0,
                    "dpmo_combined_board": 800.0,
                    "board_n": 20,
                    "combined_defect_event_count": 6,
                },
            },
            "per_measure": {"Volume": {"n": 120, "defect": {"ppm_total": 200.0, "dpmo_feature": 66666.7, "zbench_st": 2.9}}},
            "relation": {"corr_vol_area": 0.92, "corr_vol_height": 0.45, "corr_area_height": 0.42},
        },
        diagnostics=[],
        analysis_payload={
            "spatial": {"statistics": {"points": 1180}, "metadata": {"mode": "value"}},
            "pareto": {"components": [{"component_id": "R101", "abnormal_rate": 0.23}]},
        },
        report_context={
            "relation_meta": {"match_rate": 98.2, "match_count": 1180, "unmatch_count": 22},
            "filter_context": {"line": "L3", "batch": "B01", "refdes": "全部 (All)"},
            "coordinate_registry_entry": {"product_part_no": "PN-7788"},
            "product_spec_profile": {
                "product_name": "X3000",
                "stencil_type": "normal",
                "thickness_main": 0.09,
                "updated_at": "2026-03-16T12:12:29",
            },
        },
        output_path=str(output_path),
    )

    assert ok, err
    prs = Presentation(str(output_path))
    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    assert "PN-7788" in all_text
    assert "產品料號已由座標註冊表自動回填" in all_text
    assert "座標關聯成功率：98.2%" in all_text
    assert "Main Thickness: 0.090 mm" in all_text
    assert "Combined DPMO(Event)" in all_text
    assert "UNKNOWN (VERIFY)" in all_text


def test_build_pptx_report_uses_shared_risk_assessment_context(tmp_path) -> None:
    output_path = tmp_path / "builder-shared-risk-assessment.pptx"
    ok, err = builder.build_pptx_report(
        wo_master={"work_order_no": "WO-RISK-01", "product_name": "Risk Align"},
        wo_spec={},
        summary_data={
            "process": {
                "min_cpk": 1.55,
                "min_cpk_measure": "Volume",
                "verdict": "可接受",
            }
        },
        diagnostics=[],
        report_context={
            "risk_assessment": {
                "level": "HIGH",
                "level_display": "高風險 (High)",
                "error_count": 0,
                "warning_count": 0,
                "total_count": 0,
            }
        },
        output_path=str(output_path),
    )

    assert ok, err
    prs = Presentation(str(output_path))
    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    assert "製程風險 / Process Risk" in all_text
    assert "整體風險：" in all_text
    assert "高風險 (High)" in all_text


def test_build_pptx_report_fallback_risk_uses_shared_policy(tmp_path) -> None:
    output_path = tmp_path / "builder-fallback-risk-shared-policy.pptx"
    ok, err = builder.build_pptx_report(
        wo_master={"work_order_no": "WO-RISK-02", "product_name": "Risk Fallback"},
        wo_spec={},
        summary_data={
            "process": {
                "min_cpk": 1.62,
                "min_cpk_measure": "Volume",
                "verdict": "不可接受",
            }
        },
        diagnostics=[{"severity": "info"}],
        report_context={},
        output_path=str(output_path),
    )

    assert ok, err
    prs = Presentation(str(output_path))
    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    assert "製程風險 / Process Risk" in all_text
    assert "整體風險：" in all_text
    assert "高風險 (High)" in all_text


def test_build_pptx_report_fallback_risk_shows_high_priority_signal_count(tmp_path) -> None:
    output_path = tmp_path / "builder-fallback-risk-high-priority.pptx"
    ok, err = builder.build_pptx_report(
        wo_master={"work_order_no": "WO-RISK-03", "product_name": "Risk Priority"},
        wo_spec={},
        summary_data={
            "process": {
                "min_cpk": 1.55,
                "min_cpk_measure": "Volume",
                "verdict": "可接受",
            }
        },
        diagnostics=[
            {"summary": "rule-a", "severity": "info", "priority": "high"},
            {"summary": "rule-b", "severity": "info", "priority": "high"},
        ],
        report_context={},
        output_path=str(output_path),
    )

    assert ok, err
    prs = Presentation(str(output_path))
    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    assert "整體風險：" in all_text
    assert "高風險 (High)" in all_text
    assert "高優先訊號：2" in all_text


def test_build_pptx_report_adds_chart_evidence_gallery_when_images_available(
    monkeypatch,
    tmp_path,
) -> None:
    sample_png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9sL9sAAAAASUVORK5CYII="
    )

    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: sample_png,
    )

    output_path = tmp_path / "builder-chart-gallery.pptx"
    ok, err = builder.build_pptx_report(
        wo_master={"work_order_no": "WO-GALLERY-01"},
        wo_spec={},
        summary_data={},
        diagnostics=[],
        analysis_payload={"selected_features": ["Volume"]},
        report_context={
            "selected_features": ["Volume"],
            "available_features": ["Volume"],
            "evidence_coverage": build_chart_evidence_coverage(
                selected_chart_ids=["imr", "run_chart", "histogram_spec", "ooc_analysis"],
                selected_features=["Volume"],
                available_features=["Volume"],
                has_coordinate_data=True,
            ),
        },
        chart_ids_to_export=["imr", "run_chart", "histogram_spec", "ooc_analysis"],
        output_path=str(output_path),
    )

    assert ok, err
    prs = Presentation(str(output_path))
    assert len(prs.slides) >= 13
    deck_text = _deck_text(prs)
    assert "5A. Chart Evidence Gallery" in deck_text
    assert "個別值與移動極差圖" in deck_text
    assert "圖表證據覆蓋表" in deck_text
    assert "已輸出" in deck_text


def test_build_pptx_report_excludes_spatial_evidence_without_coordinate_scope(
    monkeypatch,
    tmp_path,
) -> None:
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "builder-spatial-excluded.pptx"
    report_context = {
        "selected_features": ["Volume"],
        "available_features": ["Volume"],
        "relation_meta": {"match_rate": 99.4, "match_count": 570, "unmatch_count": 2},
        "data_scope": {
            "has_coordinate_data": False,
            "sample_n": 6,
            "selected_features": ["Volume"],
            "used_sources": ["量測資料", "管制規格", "工單資訊"],
            "excluded_evidence": [
                {
                    "id": "spatial",
                    "label": "座標值 / 空間映射 / 空間熱圖",
                    "reason": "本批資料未提供有效 X/Y 座標欄位",
                }
            ],
            "section_trust": {
                "statistics": "可信：資料直接計算",
                "charts": "可信：圖表證據",
                "inference": "需複核：規則推論",
                "spatial": "未納入：資料缺失",
            },
        },
        "evidence_coverage": build_chart_evidence_coverage(
            selected_chart_ids=["spatial_heatmap", "histogram_spec"],
            selected_features=["Volume"],
            available_features=["Volume"],
            has_coordinate_data=False,
        ),
    }

    ok, err = builder.build_pptx_report(
        wo_master={"work_order_no": "WO-NO-XY-01", "product_name": "No XY"},
        wo_spec={},
        summary_data={},
        diagnostics=[],
        analysis_payload={
            "selected_features": ["Volume"],
            "spatial": {"statistics": {"points": 570}, "metadata": {"mode": "value"}},
        },
        report_context=report_context,
        chart_ids_to_export=["spatial_heatmap", "histogram_spec"],
        output_path=str(output_path),
    )

    assert ok, err
    prs = Presentation(str(output_path))
    deck_text = _deck_text(prs)
    assert "本批資料未提供座標欄位，空間分析未納入判讀。" in deck_text
    assert "圖表證據覆蓋表" in deck_text
    assert "空間熱圖" in deck_text
    assert "缺座標資料" in deck_text
    assert "證據類型：未納入" in deck_text
    assert "座標關聯成功率：99.4%" not in deck_text
    assert "空間有效點數" not in deck_text


def test_build_pptx_report_embeds_dashboard_layers_when_present(tmp_path) -> None:
    """儀表板 dashboard_layers 應出現在首頁與能力頁等（打散對齊章節）。"""
    output_path = tmp_path / "builder-dashboard-layers.pptx"
    ok, err = builder.build_pptx_report(
        wo_master={"work_order_no": "WO-DASH-01", "product_name": "X"},
        wo_spec={},
        summary_data={
            "process": {
                "verdict": "可接受",
                "min_cpk": 1.5,
                "min_cpk_measure": "Height",
                "overall_yield_pct": 99.0,
                "dashboard_layers": {
                    "layer_1_alarm": {
                        "ooc_rate": 0.01,
                        "ooc_rate_state": "Warning",
                        "anomaly_cluster_count": 2,
                    },
                    "layer_2_kpi": {"yield_pct": 99.0, "dpmo": 100.0, "sigma_level": 3.0},
                    "layer_3_info": {"driver_feature": "Height", "sample_size": 100, "range": 10.0},
                    "layer_4_defect_structure": {
                        "defect_pattern_zh": "同元件集中",
                        "cluster_ratio": 0.05,
                        "top_oos_refdes": [{"id": "R1", "oos_count": 3}],
                    },
                    "layer_5_spec_analysis": {
                        "cpk": 1.5,
                        "usl": 150.0,
                        "lsl": 70.0,
                        "target": 100.0,
                        "spec_tightness_level": "high_capability",
                        "spec_range": 80.0,
                    },
                    "layer_6_product_context": {
                        "product_name": "P1",
                        "work_order_no": "WO1",
                        "stencil_type": "S1",
                        "stencil_thickness": 0.1,
                    },
                    "layer_7_engineering_info": {"mean": 100.0, "std": 1.0},
                    "layer_8_diagnosis": {
                        "priority": "low",
                        "issue_type_display_zh": "局部",
                        "root_cause_zh": "群集傾向",
                        "recommended_action_zh": "複核焊墊",
                    },
                },
            },
            "per_measure": {},
        },
        diagnostics=[
            {
                "summary": "測試",
                "severity": "info",
                "chart_title": "t",
                "chart_bytes": None,
                "observable_charts": [],
                "evidence_lines": [],
                "ipc_lines": [],
                "recommended_actions": [],
            }
        ],
        report_context={
            "risk_assessment": {"level": "LOW"},
            "relation_meta": {"match_rate": 99.0},
            "selected_features": ["Volume"],
            "available_features": ["Volume"],
        },
        analysis_payload={"selected_features": ["Volume"]},
        output_path=str(output_path),
    )
    assert ok, err
    prs = Presentation(str(output_path))
    deck_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    assert "儀表板" in deck_text or "OOC" in deck_text
    assert "製程能力" in deck_text
