import pandas as pd
from pptx import Presentation
from pptx.util import Mm

from app.data.session_store import SessionStore, _analysis_cache_key
from app.services.report_service import (
    ReportService,
    _build_executive_summary_html,
    _build_pptx_diagnostics,
    _build_report_html,
    _collect_pptx_actions,
    _compute_risk_level,
    _display_name_to_chart_id,
    _format_pptx_evidence_lines,
    _get_pptx_chart_title,
    _make_cached_chart_renderer,
    _normalize_pptx_observable_charts,
    _normalize_pptx_severity,
    _pptx_severity_rank,
)


def _all_slide_text(prs: Presentation) -> str:
    """Collect full deck text for content assertions independent of slide index."""
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            parts.append(cell.text)
            elif hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def test_display_name_to_chart_id_handles_hint_aliases() -> None:
    assert _display_name_to_chart_id("CUSUM 圖") == "cusum"
    assert _display_name_to_chart_id("製程能力圖 (Capability)") == "histogram_spec"
    assert _display_name_to_chart_id("空間熱圖 (Spatial Heatmap)") == "spatial_heatmap"


def test_collect_pptx_actions_prefers_rule_specific_failure_mode_actions() -> None:
    drying_actions = _collect_pptx_actions({}, rule_id="volume_decline_along_board", limit=3)
    cpk_actions = _collect_pptx_actions({}, rule_id="cpk_below_threshold", limit=3)

    assert "縮短開罐後使用時間" in drying_actions
    assert "確認設備參數穩定性" in cpk_actions
    assert "縮短開罐後使用時間" not in cpk_actions


def test_format_pptx_evidence_lines_formats_ratio_units_correctly() -> None:
    lines = _format_pptx_evidence_lines(
        {
            "variance_ratio": 2.1,
            "ooc_ratio": 0.125,
            "cv": 0.2,
        },
        limit=4,
    )

    assert "Variance Ratio: 2.10x" in lines
    assert "OOC Ratio: 12.5%" in lines
    assert "CV: 20.0%" in lines


def test_build_pptx_diagnostics_assigns_rule_specific_actions(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )
    payload = {
        "run_chart": {
            "data": {"values": [111, 110, 109, 108, 107, 106, 105, 104, 103, 102, 101, 100]},
        },
        "cap": {
            "metadata": {"is_valid": True, "usl": 110.0, "lsl": 90.0},
            "statistics": {"cpk": 0.82, "cp": 1.01, "mean": 99.5, "sigma_st": 4.8},
        },
    }

    diagnostics = _build_pptx_diagnostics(payload, ["Volume"])
    by_rule = {item["rule_id"]: item["recommended_actions"] for item in diagnostics}

    assert "volume_decline_along_board" in by_rule
    assert "cpk_below_threshold" in by_rule
    assert "縮短開罐後使用時間" in by_rule["volume_decline_along_board"]
    assert "確認設備參數穩定性" in by_rule["cpk_below_threshold"]
    assert "縮短開罐後使用時間" not in by_rule["cpk_below_threshold"]
    assert diagnostics[0]["evidence_type"] == "統計計算 / 規則推論"


def test_build_pptx_diagnostics_expands_multi_feature_parameters(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )
    payload = {
        "selected_features": ["Volume", "Height"],
        "parameters": {
            "Volume": {
                "cap": {
                    "metadata": {"is_valid": True, "usl": 110.0, "lsl": 90.0},
                    "statistics": {"cpk": 0.82, "cp": 1.01, "mean": 99.5, "sigma_st": 4.8},
                }
            },
            "Height": {
                "cap": {
                    "metadata": {"is_valid": True, "usl": 0.18, "lsl": 0.12},
                    "statistics": {"cpk": 1.67, "cp": 1.80, "mean": 0.15, "sigma_st": 0.004},
                }
            },
        },
    }

    diagnostics = _build_pptx_diagnostics(payload, ["Volume", "Height"])

    assert len(diagnostics) == 1
    assert diagnostics[0]["feature_label"] == "Volume (體積)"
    assert "[Volume (體積)]" in diagnostics[0]["summary"]
    assert diagnostics[0]["chart_title"].startswith("Volume (體積) |")


def test_build_pptx_diagnostics_keeps_first_relevant_chart_title_when_render_missing(
    monkeypatch,
) -> None:
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )
    payload = {
        "cap": {
            "metadata": {"is_valid": True, "usl": 110.0, "lsl": 90.0},
            "statistics": {"cpk": 0.82, "cp": 1.01, "mean": 99.5, "sigma_st": 4.8},
        }
    }

    diagnostics = _build_pptx_diagnostics(payload, ["Volume"])

    expected_prefix = f"Volume (體積) | {_get_pptx_chart_title('histogram_spec')}"
    assert diagnostics[0]["chart_title"] == expected_prefix


def test_build_pptx_diagnostics_includes_missing_chart_reason(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )
    payload = {
        "spatial": {
            "metadata": {
                "is_valid": False,
                "error": "缺乏有效座標映射資料",
            },
            "data": {
                "x": [0, 10, 20],
                "y": [0, 5, 10],
            },
            "modes": {
                "oos_density": {
                    "values": [5, 0, 5],
                }
            },
        }
    }

    diagnostics = _build_pptx_diagnostics(
        payload,
        ["Volume"],
    )

    spatial_item = next(
        item for item in diagnostics if item["rule_id"] == "edge_spatial_cluster"
    )
    assert "缺乏有效座標映射資料" in spatial_item["chart_missing_reason"]


def test_normalize_pptx_observable_charts_deduplicates_verbose_aliases() -> None:
    titles = _normalize_pptx_observable_charts(
        ["製程能力圖 (Capability)", "直方圖 (Histogram)", "常態機率圖 (Normality)"]
    )

    assert titles == ["分布與能力", "常態分析"]


def test_build_pptx_diagnostics_sorts_high_severity_before_warnings(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )
    payload = {
        "selected_features": ["Volume", "Height"],
        "parameters": {
            "Volume": {
                "normality": {
                    "metadata": {"is_valid": True},
                    "statistics": {
                        "p_value": 0.01,
                        "is_normal": False,
                        "total_n": 12,
                        "test_name": "Shapiro-Wilk",
                        "r_squared": 0.92,
                    },
                }
            },
            "Height": {
                "cap": {
                    "metadata": {"is_valid": True, "usl": 0.18, "lsl": 0.12},
                    "statistics": {"cpk": 0.79, "cp": 0.95, "mean": 0.151, "sigma_st": 0.012},
                }
            },
        },
    }

    diagnostics = _build_pptx_diagnostics(payload, ["Volume", "Height"])

    assert [item["severity"] for item in diagnostics] == ["error", "warning"]
    assert diagnostics[0]["feature_label"] == "Height (高度)"
    assert diagnostics[1]["feature_label"] == "Volume (體積)"


def test_build_pptx_diagnostics_ignores_info_only_hints(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.analytics.root_cause_engine.infer_root_cause_hints",
        lambda *_args, **_kwargs: [
            {
                "hint": "資料穩定，僅供觀察",
                "rule_id": "info_only",
                "severity": "info",
                "observable_charts": ["趨勢圖"],
            }
        ],
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    diagnostics = _build_pptx_diagnostics({"run_chart": {"data": {"values": [1, 2, 3]}}}, ["Volume"])

    assert diagnostics == []


def test_build_pptx_diagnostics_normalizes_uppercase_severity(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.analytics.root_cause_engine.infer_root_cause_hints",
        lambda *_args, **_kwargs: [
            {
                "hint": "Cpk 低於門檻",
                "rule_id": "cpk_below_threshold",
                "severity": "ERROR",
                "observable_charts": ["製程能力圖 (Capability)"],
            }
        ],
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    diagnostics = _build_pptx_diagnostics({"cap": {"statistics": {"cpk": 0.82}}}, ["Volume"])

    assert len(diagnostics) == 1
    assert diagnostics[0]["severity"] == "error"


def test_build_pptx_diagnostics_skips_empty_hint_text(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.analytics.root_cause_engine.infer_root_cause_hints",
        lambda *_args, **_kwargs: [
            {
                "hint": "   ",
                "rule_id": "empty_hint",
                "severity": "warning",
                "observable_charts": ["趨勢圖"],
            },
            {
                "hint": "Volume 端出現持續下降趨勢",
                "rule_id": "volume_decline_along_board",
                "severity": "warning",
                "observable_charts": ["趨勢圖"],
            },
        ],
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    diagnostics = _build_pptx_diagnostics({"run_chart": {"data": {"values": [1, 2, 3]}}}, ["Volume"])

    assert len(diagnostics) == 1
    assert diagnostics[0]["rule_id"] == "volume_decline_along_board"


def test_build_pptx_diagnostics_supports_risk_only_mode_without_chart_render(
    monkeypatch,
) -> None:
    monkeypatch.setattr(
        "app.analytics.root_cause_engine.infer_root_cause_hints",
        lambda *_args, **_kwargs: [
            {
                "hint": "CUSUM 偏移告警",
                "rule_id": "cusum_trend_drift",
                "severity": "warning",
                "observable_charts": ["CUSUM 圖"],
            }
        ],
    )

    diagnostics = _build_pptx_diagnostics(
        {"cusum": {"data": {"values": [1, 2, 3]}}},
        ["Volume"],
        include_chart_render=False,
    )

    assert len(diagnostics) == 1
    assert diagnostics[0]["chart_bytes"] is None
    assert diagnostics[0]["chart_missing_reason"] == "風險彙總模式：略過圖像渲染。"


def test_build_pptx_diagnostics_preserves_priority_for_risk_scoring(
    monkeypatch,
) -> None:
    monkeypatch.setattr(
        "app.analytics.root_cause_engine.infer_root_cause_hints",
        lambda *_args, **_kwargs: [
            {
                "hint": "Edge cluster warning",
                "rule_id": "edge_spatial_cluster",
                "severity": "warning",
                "priority": "high",
                "observable_charts": ["空間熱圖"],
            }
        ],
    )
    diagnostics = _build_pptx_diagnostics(
        {"spatial": {"data": {"x": [1], "y": [1]}}},
        ["Volume"],
        include_chart_render=False,
    )
    assert len(diagnostics) == 1
    assert diagnostics[0]["priority"] == "high"


def test_pptx_severity_rank_orders_error_warning_info() -> None:
    assert _pptx_severity_rank("error") < _pptx_severity_rank("warning")
    assert _pptx_severity_rank("warning") < _pptx_severity_rank("info")


def test_normalize_pptx_severity_supports_aliases_and_priority() -> None:
    assert _normalize_pptx_severity("HIGH") == "error"
    assert _normalize_pptx_severity("medium") == "warning"
    assert _normalize_pptx_severity("warn") == "warning"
    assert _normalize_pptx_severity("unknown", priority="high") == "error"
    assert _normalize_pptx_severity(None) == "info"


def test_compute_risk_level_respects_unacceptable_verdict_floor() -> None:
    level = _compute_risk_level([], process={"verdict": "不可接受"})
    assert level == "HIGH"


def test_compute_risk_level_normalizes_severity_aliases_in_hints() -> None:
    level = _compute_risk_level([{"severity": "WARN"}, {"severity": "HIGH"}], process={"verdict": "可接受"})
    assert level == "HIGH"


def test_compute_risk_level_uses_diagnostics_signal_when_provided() -> None:
    level = _compute_risk_level(
        [{"severity": "info"}],
        process={"verdict": "可接受"},
        diagnostics=[{"severity": "warning"}],
    )
    assert level == "MEDIUM"


def test_executive_summary_risk_badge_uses_diagnostics_signals() -> None:
    html = _build_executive_summary_html(
        hints=[],
        ro_payload={},
        summary_data={"process": {"verdict": "可接受"}},
        total_n=120,
        batch_qty=20,
        diagnostics=[{"severity": "warning"}],
    )
    assert "MEDIUM 中風險" in html


def test_build_report_html_includes_shared_risk_snapshot() -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True}
    store.coord_meta = {"is_valid": True}
    store.relation_meta = {"match_rate": 95.0}
    store.selected_features = ["Volume"]
    store.last_analysis_payload = {
        "selected_features": ["Volume"],
        "summary": {
            "process": {
                "verdict": "不可接受",
            }
        },
        "run_chart": {
            "data": {"values": [111, 110, 109, 108, 107, 106, 105, 104, 103, 102, 101, 100]},
        },
    }

    html = _build_report_html(store, chart_ids_to_export=[], selected_features=["Volume"])

    assert "[0] 風險摘要 (Shared Decision)" in html
    assert "Risk Level: 高風險 (High)" in html
    assert "HighPriority=" in html


def test_build_report_html_tolerates_non_numeric_rate_fields() -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True}
    store.coord_meta = {"is_valid": True}
    store.relation_meta = {"match_rate": "N/A"}
    store.selected_features = ["Volume"]
    store.last_analysis_payload = {
        "pareto": {
            "components": [
                {"component_id": "U1", "abnormal_rate": "bad-value", "total": 12},
            ]
        }
    }

    html = _build_report_html(store, chart_ids_to_export=[], selected_features=["Volume"])

    assert "空間映射成功率: 0.0%" in html
    assert "U1: abnormal_rate=0.00%, n=12" in html


def test_build_pptx_diagnostics_skips_non_list_hint_payload(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.analytics.root_cause_engine.infer_root_cause_hints",
        lambda *_args, **_kwargs: {"hint": "unexpected"},
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )
    diagnostics = _build_pptx_diagnostics({"cap": {"statistics": {"cpk": 0.82}}}, ["Volume"])
    assert diagnostics == []


def test_build_pptx_diagnostics_maps_high_alias_to_error(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.analytics.root_cause_engine.infer_root_cause_hints",
        lambda *_args, **_kwargs: [
            {
                "hint": "製程能力急遽惡化",
                "rule_id": "cpk_below_threshold",
                "severity": "HIGH",
                "priority": "high",
                "observable_charts": ["製程能力圖 (Capability)"],
            }
        ],
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )
    diagnostics = _build_pptx_diagnostics({"cap": {"statistics": {"cpk": 0.82}}}, ["Volume"])
    assert len(diagnostics) == 1
    assert diagnostics[0]["severity"] == "error"


def test_generate_pptx_report_creates_a4_landscape_deck(monkeypatch, tmp_path) -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True, "missing_required": []}
    store.coord_meta = {"is_valid": True, "missing_required": []}
    store.relation_meta = {"match_rate": 98.2, "unmatch_count": 1}
    store.workorder_master = {
        "work_order_no": "WO-7788",
        "product_name": "Demo Board",
        "product_part_no": "P-1001",
        "supplier": "ACME",
        "batch_qty": 120,
    }
    store.workorder_spec = {
        "volume": {"usl": "110", "lsl": "90", "target": "100"},
        "area": {"usl": "105", "lsl": "95", "target": "100"},
        "height": {"usl": "0.18", "lsl": "0.12", "target": "0.15"},
    }
    store.selected_features = ["Volume"]
    store.meas_df = pd.DataFrame(
        {
            "BoardNo": [f"B{i:02d}" for i in range(1, 13)],
            "RefDes": ["R1"] * 12,
            "Volume": [110, 108, 107, 105, 104, 103, 100, 99, 98, 96, 94, 92],
            "Area": [101, 100, 102, 101, 100, 99, 101, 100, 99, 98, 97, 96],
            "Height": [0.152, 0.151, 0.150, 0.149, 0.149, 0.148, 0.147, 0.146, 0.145, 0.145, 0.144, 0.143],
        }
    )
    store.last_analysis_payload = {
        "run_chart": {
            "data": {
                "values": [110, 108, 107, 105, 104, 103, 100, 99, 98, 96, 94, 92],
            }
        }
    }

    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "spi-analysis-report.pptx"
    ok, err = ReportService().generate_pptx_report(str(output_path))

    assert ok, err
    assert output_path.exists()

    prs = Presentation(str(output_path))
    assert prs.slide_width == Mm(297)
    assert prs.slide_height == Mm(210)
    assert len(prs.slides) >= 12

    def _slide_text(slide) -> str:
        return "\n".join(
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )

    slide_texts = [_slide_text(slide) for slide in prs.slides]
    combined = "\n".join(slide_texts)
    assert "1. Product & Work Order Information" in combined
    assert "2. Control Specification" in combined
    assert "3. Statistics Summary" in combined
    assert "異常診斷詳頁" in combined or "Anomaly Diagnosis" in combined
    assert "建議動作" in combined or "建議" in combined

    stats_slide = next(
        (s for s in prs.slides if "3. Statistics Summary" in _slide_text(s)),
        None,
    )
    assert stats_slide is not None
    stats_tables = [
        shape.table for shape in stats_slide.shapes if getattr(shape, "has_table", False)
    ]
    assert stats_tables
    stats_table = stats_tables[0]
    headers = [stats_table.cell(0, ci).text for ci in range(len(stats_table.columns))]
    volume_col = headers.index("Volume (體積)")
    yield_row = None
    for ri in range(len(stats_table.rows)):
        if stats_table.cell(ri, 0).text == "Yield (良率 %)":
            yield_row = ri
            break
    assert yield_row is not None
    assert "%" in stats_table.cell(yield_row, volume_col).text


def test_generate_pptx_report_shows_missing_chart_reason_in_placeholder(
    monkeypatch,
    tmp_path,
) -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True, "missing_required": []}
    store.workorder_master = {"work_order_no": "WO-SPATIAL-2201", "product_name": "Spatial Missing"}
    store.workorder_spec = {
        "volume": {"usl": "110", "lsl": "90", "target": "100"},
    }
    store.selected_features = ["Volume"]
    store.meas_df = pd.DataFrame(
        {
            "BoardNo": [f"B{i:02d}" for i in range(1, 7)],
            "RefDes": ["R1"] * 6,
            "Volume": [100.0] * 6,
        }
    )

    def _fake_compute_analysis_payload(
        filtered_df: pd.DataFrame,
        selected_features: list[str],
        usl: float,
        lsl: float,
        target: float,
        workorder_spec: dict[str, object] | None = None,
        workorder_master: dict[str, object] | None = None,
        cancel_fn=None,
    ) -> tuple[dict[str, object], None]:
        return (
            {
                "selected_features": selected_features,
                "spatial": {
                    "metadata": {
                        "is_valid": False,
                        "error": "缺乏有效座標映射資料",
                    },
                    "data": {
                        "x": [0, 10, 20],
                        "y": [0, 5, 10],
                    },
                    "modes": {
                        "oos_density": {
                            "values": [5, 0, 5],
                        }
                    },
                },
            },
            None,
        )

    monkeypatch.setattr(
        "app.viewmodels.chart_analysis_viewmodel.compute_analysis_payload",
        _fake_compute_analysis_payload,
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "missing-chart-reason-report.pptx"
    ok, err = ReportService().generate_pptx_report(str(output_path))

    assert ok, err
    prs = Presentation(str(output_path))
    deck_text = _all_slide_text(prs)
    assert "缺乏有效座標映射資料" in deck_text


def test_generate_pptx_report_marks_spatial_unincluded_without_xy(
    monkeypatch,
    tmp_path,
) -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True, "missing_required": []}
    store.relation_meta = {"match_rate": 99.4, "match_count": 6, "unmatch_count": 0}
    store.workorder_master = {"work_order_no": "WO-NO-XY-03", "product_name": "No XY Lot"}
    store.workorder_spec = {"volume": {"usl": "110", "lsl": "90", "target": "100"}}
    store.selected_features = ["Volume"]
    store.meas_df = pd.DataFrame(
        {
            "BoardNo": [f"B{i:02d}" for i in range(1, 7)],
            "RefDes": ["R1"] * 6,
            "Volume": [100.0, 99.8, 100.2, 100.1, 99.9, 100.0],
        }
    )

    def _fake_compute_analysis_payload(
        filtered_df: pd.DataFrame,
        selected_features: list[str],
        usl: float,
        lsl: float,
        target: float,
        workorder_spec: dict[str, object] | None = None,
        workorder_master: dict[str, object] | None = None,
        cancel_fn=None,
    ) -> tuple[dict[str, object], None]:
        return (
            {
                "selected_features": selected_features,
                "spatial": {
                    "metadata": {"is_valid": True, "mode": "value"},
                    "statistics": {"points": 6},
                },
            },
            None,
        )

    monkeypatch.setattr(
        "app.viewmodels.chart_analysis_viewmodel.compute_analysis_payload",
        _fake_compute_analysis_payload,
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "service-no-coordinate-report.pptx"
    ok, err = ReportService().generate_pptx_report(
        str(output_path),
        chart_ids_to_export=["spatial_heatmap", "histogram_spec"],
    )

    assert ok, err
    prs = Presentation(str(output_path))
    deck_text = _all_slide_text(prs)
    assert "本批資料未提供座標欄位，空間分析未納入判讀。" in deck_text
    assert "圖表證據覆蓋表" in deck_text
    assert "缺座標資料" in deck_text
    assert "座標關聯成功率：99.4%" not in deck_text
    assert "空間有效點數" not in deck_text


def test_generate_pptx_report_keeps_multi_feature_diagnostics_visible(
    monkeypatch,
    tmp_path,
) -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True, "missing_required": []}
    store.workorder_master = {"work_order_no": "WO-8801", "product_name": "Multi Feature Lot"}
    store.workorder_spec = {
        "volume": {"usl": "110", "lsl": "90", "target": "100"},
        "height": {"usl": "0.18", "lsl": "0.12", "target": "0.15"},
    }
    store.selected_features = ["Volume", "Height"]
    store.meas_df = pd.DataFrame(
        {
            "BoardNo": [f"B{i:02d}" for i in range(1, 7)],
            "RefDes": ["R1"] * 6,
            "Volume": [109.0, 107.0, 105.0, 103.0, 101.0, 99.0],
            "Height": [0.150, 0.151, 0.149, 0.150, 0.151, 0.150],
        }
    )

    def _fake_compute_analysis_payload(
        filtered_df: pd.DataFrame,
        selected_features: list[str],
        usl: float,
        lsl: float,
        target: float,
        workorder_spec: dict[str, object] | None = None,
        workorder_master: dict[str, object] | None = None,
        cancel_fn=None,
    ) -> tuple[dict[str, object], None]:
        return (
            {
                "selected_features": selected_features,
                "parameters": {
                    "Volume": {
                        "cap": {
                            "metadata": {"is_valid": True, "usl": 110.0, "lsl": 90.0},
                            "statistics": {"cpk": 0.82, "cp": 1.01, "mean": 99.5, "sigma_st": 4.8},
                        }
                    },
                    "Height": {
                        "cap": {
                            "metadata": {"is_valid": True, "usl": 0.18, "lsl": 0.12},
                            "statistics": {"cpk": 1.67, "cp": 1.80, "mean": 0.15, "sigma_st": 0.004},
                        }
                    },
                },
            },
            None,
        )

    monkeypatch.setattr(
        "app.viewmodels.chart_analysis_viewmodel.compute_analysis_payload",
        _fake_compute_analysis_payload,
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "multi-feature-report.pptx"
    ok, err = ReportService().generate_pptx_report(str(output_path))

    assert ok, err
    prs = Presentation(str(output_path))
    assert len(prs.slides) >= 12
    deck_text = _all_slide_text(prs)
    assert "Volume (體積)" in deck_text
    assert "確認設備參數穩定性" in deck_text
    assert "未觸發異常診斷規則" not in deck_text


def test_generate_pptx_report_prefers_fresh_filtered_payload_over_stale_cache(
    monkeypatch,
    tmp_path,
) -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True, "missing_required": []}
    store.workorder_master = {"work_order_no": "WO-2201", "product_name": "Filtered Lot"}
    store.workorder_spec = {
        "volume": {"usl": "110", "lsl": "90", "target": "100"},
    }
    store.selected_features = ["Volume"]
    store.filter_batch = "B01"
    store.meas_df = pd.DataFrame(
        {
            "BoardNo": ["B01", "B01", "B02", "B02"],
            "RefDes": ["R1", "R1", "R1", "R1"],
            "Volume": [100.0, 99.8, 108.0, 107.5],
        }
    )
    store.last_analysis_payload = {
        "run_chart": {
            "data": {
                "values": [111, 110, 109, 108, 107, 106, 105, 104, 103, 102, 101, 100],
            }
        }
    }

    captured: dict[str, object] = {}

    def _fake_compute_analysis_payload(
        filtered_df: pd.DataFrame,
        selected_features: list[str],
        usl: float,
        lsl: float,
        target: float,
        workorder_spec: dict[str, object] | None = None,
        workorder_master: dict[str, object] | None = None,
        cancel_fn=None,
    ) -> tuple[dict[str, object], None]:
        captured["rows"] = len(filtered_df)
        captured["selected_features"] = selected_features
        return (
            {
                "cap": {
                    "metadata": {"is_valid": True, "usl": usl, "lsl": lsl},
                    "statistics": {"cpk": 0.82, "cp": 1.01, "mean": 99.5, "sigma_st": 4.8},
                }
            },
            None,
        )

    monkeypatch.setattr(
        "app.viewmodels.chart_analysis_viewmodel.compute_analysis_payload",
        _fake_compute_analysis_payload,
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "fresh-payload-report.pptx"
    ok, err = ReportService().generate_pptx_report(str(output_path))

    assert ok, err
    assert captured["rows"] == 2
    assert captured["selected_features"] == ["Volume"]

    prs = Presentation(str(output_path))
    deck_text = _all_slide_text(prs)
    assert "確認設備參數穩定性" in deck_text
    assert "縮短開罐後使用時間" not in deck_text


def test_generate_pptx_report_reuses_matching_analysis_cache_without_recompute(
    monkeypatch,
    tmp_path,
) -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True, "missing_required": []}
    store.workorder_master = {"work_order_no": "WO-CACHE", "product_name": "Cached Lot"}
    store.workorder_spec = {
        "volume": {"usl": "110", "lsl": "90", "target": "100"},
    }
    store.selected_features = ["Volume"]
    store.filter_batch = "B01"
    store.meas_df = pd.DataFrame(
        {
            "BoardNo": ["B01", "B01", "B02", "B02"],
            "RefDes": ["R1", "R1", "R1", "R1"],
            "Volume": [100.0, 99.8, 108.0, 107.5],
        }
    )
    cache_key = _analysis_cache_key(
        ["Volume"],
        "B01",
        "全部 (All)",
        "全部 (All)",
        spec_version=store.spec_cache_token(store.workorder_spec),
    )
    store._analysis_cache[cache_key] = {
        "selected_features": ["Volume"],
        "_ctx_batch": "B01",
        "_ctx_refdes": "全部 (All)",
        "_ctx_part_type": "全部 (All)",
        "run_chart": {"data": {"values": [100.0, 99.8]}},
    }

    def _should_not_recompute(*args, **kwargs):  # noqa: ARG001
        raise AssertionError("cached report export should not recompute analysis payload")

    monkeypatch.setattr(
        "app.viewmodels.chart_analysis_viewmodel.compute_analysis_payload",
        _should_not_recompute,
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "cached-payload-report.pptx"
    ok, err = ReportService().generate_pptx_report(str(output_path))

    assert ok, err
    assert output_path.is_file()


def test_cached_chart_renderer_reuses_same_export_image(monkeypatch) -> None:
    calls = {"count": 0}

    def _fake_render(*args, **kwargs):  # noqa: ARG001
        calls["count"] += 1
        return b"png-bytes"

    monkeypatch.setattr("app.services.chart_render.render_chart_to_png_bytes", _fake_render)
    stats = {"requests": 0, "hits": 0, "misses": 0}
    render = _make_cached_chart_renderer(stats)

    first = render("imr", {"selected_features": ["Volume"]}, features=["Volume"], context="report")
    second = render("imr", {"selected_features": ["Volume"]}, features=["Volume"], context="report")

    assert first == second == b"png-bytes"
    assert calls["count"] == 1
    assert stats == {"requests": 2, "hits": 1, "misses": 1}


def test_generate_pptx_report_falls_back_to_cached_payload_when_recompute_fails(
    monkeypatch,
    tmp_path,
) -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True, "missing_required": []}
    store.workorder_master = {"work_order_no": "WO-2202", "product_name": "Fallback Lot"}
    store.workorder_spec = {
        "volume": {"usl": "110", "lsl": "90", "target": "100"},
    }
    store.selected_features = ["Volume"]
    store.meas_df = pd.DataFrame(
        {
            "BoardNo": [f"B{i:02d}" for i in range(1, 13)],
            "RefDes": ["R1"] * 12,
            "Volume": [100.0] * 12,
        }
    )
    store.last_analysis_payload = {
        "run_chart": {
            "data": {
                "values": [111, 110, 109, 108, 107, 106, 105, 104, 103, 102, 101, 100],
            }
        }
    }

    monkeypatch.setattr(
        "app.viewmodels.chart_analysis_viewmodel.compute_analysis_payload",
        lambda *args, **kwargs: (None, "forced failure"),
    )
    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "fallback-payload-report.pptx"
    ok, err = ReportService().generate_pptx_report(str(output_path))

    assert ok, err

    prs = Presentation(str(output_path))
    deck_text = _all_slide_text(prs)
    assert "縮短開罐後使用時間" in deck_text


def test_generate_pptx_report_ignores_stale_cache_when_no_analyzable_features(
    monkeypatch,
    tmp_path,
) -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True, "missing_required": []}
    store.workorder_master = {"work_order_no": "WO-NO-FEAT", "product_name": "No Feature Columns"}
    store.workorder_spec = {
        "volume": {"usl": "110", "lsl": "90", "target": "100"},
    }
    store.selected_features = ["Volume"]
    # Deliberately omit Volume/Area/Height columns.
    store.meas_df = pd.DataFrame(
        {
            "BoardNo": [f"B{i:02d}" for i in range(1, 7)],
            "RefDes": ["R1"] * 6,
        }
    )
    # Stale cached payload that would otherwise trigger anomaly hints.
    store.last_analysis_payload = {
        "run_chart": {
            "data": {
                "values": [111, 110, 109, 108, 107, 106, 105, 104, 103, 102, 101, 100],
            }
        }
    }

    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "no-feature-columns-report.pptx"
    ok, err = ReportService().generate_pptx_report(str(output_path))

    assert ok, err
    prs = Presentation(str(output_path))
    assert len(prs.slides) >= 12
    deck_text = _all_slide_text(prs)
    assert "未觸發異常診斷規則" in deck_text
    assert "縮短開罐後使用時間" not in deck_text


def test_generate_pptx_report_keeps_fixed_twelve_sections_without_hints(monkeypatch, tmp_path) -> None:
    store = SessionStore()
    store.clear()
    store.meas_meta = {"is_valid": True, "missing_required": []}
    store.workorder_master = {"work_order_no": "WO-0001", "product_name": "Stable Lot"}
    store.workorder_spec = {
        "volume": {"usl": "110", "lsl": "90", "target": "100"},
    }
    store.selected_features = ["Volume"]
    store.meas_df = pd.DataFrame(
        {
            "BoardNo": [f"B{i:02d}" for i in range(1, 9)],
            "RefDes": ["R1"] * 8,
            "Volume": [100.0, 100.3, 99.8, 100.1, 100.2, 99.9, 100.0, 100.1],
        }
    )
    store.last_analysis_payload = {}

    monkeypatch.setattr(
        "app.services.chart_render.render_chart_to_png_bytes",
        lambda *args, **kwargs: None,
    )

    output_path = tmp_path / "stable-analysis-report.pptx"
    ok, err = ReportService().generate_pptx_report(str(output_path))

    assert ok, err
    prs = Presentation(str(output_path))
    assert len(prs.slides) >= 12
    deck_text = _all_slide_text(prs)
    assert "未觸發異常診斷規則" in deck_text
